import mimetypes
import time
import traceback
import asyncio
from contextlib import asynccontextmanager
import json
import logging
import pathlib
from datetime import datetime
from typing import Any

from langfuse import get_client, propagate_attributes
from pptx import presentation

from app.config import settings
from app.services.presentation import agent, docparser, template
from app.services.presentation.job import JobManager, JobDetail
from app.utils.storage import Storage


class PresentationService:

    def __init__(self, storage: Storage, jobs: JobManager, do_transform: bool, max_auto_retries: int, max_planner_tasks: int, max_designer_tasks: int, max_finalizer_tasks: int):
        self.storage = storage
        self.jobs = jobs
        self.do_transform = do_transform
        self.max_auto_retries = max_auto_retries
        self.planner_sem = asyncio.Semaphore(max_planner_tasks)
        self.designer_sem = asyncio.Semaphore(max_designer_tasks)
        self.finalizer_sem = asyncio.Semaphore(max_finalizer_tasks)
        self.logger = logging.getLogger(__name__)
        self.processing: dict[bytes, asyncio.Task] = {}


    @asynccontextmanager
    async def run(self):
        async with self.jobs:
            t1 = asyncio.create_task(self._run_jobs())
            t2 = asyncio.create_task(self._propagate_jobs())
            try:
                yield
            finally:
                t1.cancel()
                t2.cancel()
                await asyncio.gather(t1, t2, return_exceptions=True)


    async def _run_jobs(self):
        async for job_id in self.jobs.sub():
            if job_id is None:
                break

            job = await self.jobs.get(job_id)
            if job is None:
                if job_id.bytes in self.processing:
                    self.processing[job_id.bytes].cancel()
                    self.logger.info("Terminated job - %s", job_id)
                continue

            if job.id.bytes in self.processing or job.status == "complete":
                continue

            self.logger.info("Processing job %s - %s", job.id, job.status)
            task = asyncio.create_task(self._run_job(job))
            self.processing[job.id.bytes] = task
            task.add_done_callback(lambda t, job=job: asyncio.create_task(self._on_job_complete(job, t)))


    async def _propagate_jobs(self):
        q = asyncio.Queue()
        self.jobs.listeners.add(q)
        try:
            while True:
                id = await q.get()
                if id is None:
                    break
                job = await self.jobs.get(id)
                if job is not None:
                    await self.jobs.log(job.id, "update", json.loads(job.model_dump_json()))
        finally:
            self.jobs.listeners.discard(q)


    async def _process_checkpoint(self, job: JobDetail, event: agent.ShikshaCheckpointEvent, metadata_dir: str, prs: presentation.Presentation | None, out_path: str):
        updates = {}
        if event.metadata is not None:
            updates["metadata." + metadata_dir] = event.metadata
        if event.message is not None:
            updates["message"] = event.message

        kwargs: dict[str, Any] = {}
        if updates: kwargs["fields_set"] = updates
        if event.reason == "op": kwargs["fields_unset"] = ["metadata.error"]
        if len(kwargs) > 0: await self.jobs.update(job.id, **kwargs)

        if prs:
            await docparser.save_pptx(self.storage, prs, out_path)


    async def _on_job_complete(self, job: JobDetail, task: asyncio.Task):
        del self.processing[job.id.bytes]
        e = task.exception()
        if e:
            self.logger.info("Job raised exception - %s - %s", job.id, str(e))
            traceback.print_exception(e)
            await self.jobs.update(job.id, self._error_metadata(job, f"Unexpected error: {str(e)}", str(e)))
        elif task.cancelled():
            self.logger.info("Job cancelled - %s", job.id)
            await self.jobs.update(job.id, self._error_metadata(job, "Task was cancelled", "Task was cancelled"))


    def _error_metadata(self, job: JobDetail, message: str, error_message: str):
        attempt = job.metadata.get("error", {}).get("attempt", 0) + 1
        return {
            "status": "error",
            "message": message,
            "metadata.error.error_type": "unexpected_error",
            "metadata.error.timestamp": datetime.now().isoformat(),
            "metadata.error.error_message": error_message,
            "metadata.error.attempt": attempt,
            "metadata.error.next_attempt": time.time() + min(60, 5 * (2 ** (attempt - 1))),
            "metadata.error.attempting_recovery": True,
        }


    async def _run_job(self, job: JobDetail):
        langfuse = get_client()
        trace_id = langfuse.create_trace_id(seed=f"presentation:{job.id}")
        with langfuse.start_as_current_observation(as_type="span", name=job.status, input=job, trace_context={"trace_id": trace_id}):
            with propagate_attributes(trace_name="Shiksha-Presentation", user_id=job.user_id, session_id=f"pres-{job.id}", metadata={"tags": ",".join(job.tags)}):
                await self._actually_run_job(job)


    async def _actually_run_job(self, job: JobDetail):
        if job.status == "init":
            await self.jobs.update(job.id, {"status": "extracting_figures", "message": "Extracting figures from textbook"})
        elif job.status == "extracting_figures":
            await self.jobs.update(job.id, {"message": "Reading textbook and extracting figures"})
            stem = pathlib.Path(job.textbook_file).stem
            async with self.planner_sem, self.storage.read(self.storage.path("uploads", job.textbook_file)) as f:
                figures = await docparser.read_figures(self.storage, f, self.storage.path("out", stem, "figures"))
                await self.jobs.update(job.id, {"metadata.analysis": {
                    "extraction_time": datetime.now().isoformat(),
                    "figures": list(map(lambda x: x.model_dump(), figures)),
                }})
                if self.do_transform and ("transform_path" not in job.metadata or (job.metadata["transform_path"] is not None and not await self.storage.exists(self.storage.path("out", stem, job.metadata["transform_path"])))):
                    await self.jobs.update(job.id, {"message": "Simplifying document"})
                    # we are deliberately storing transformation in out/stem/stem instead of out/stem/jobid - if we had
                    # computed transformation for this document during another job, we can skip recomputing for this job.
                    transform_path = await docparser.transform(self.storage, f, job.textbook_mime, self.storage.path("out", stem, stem))
                    await self.jobs.update(job.id, {"metadata.transform_path": transform_path})
            await self.jobs.update(job.id, {"status": "planning_structure"})
        elif job.status == "planning_structure":
            await self.jobs.update(job.id, {"message": "Creating presentation outline"})
            stem = pathlib.Path(job.textbook_file).stem
            out_path = self.storage.path("out", stem, f"{job.id}.pptx")
            metadata = job.metadata.get("plan", {})
            if not isinstance(metadata, dict):
                metadata = {}
            if transform_path := job.metadata.get("transform_path", None):
                source_path = self.storage.path("out", stem, transform_path)
                source_mime = mimetypes.guess_type(source_path)[0]
                assert source_mime is not None
            else:
                source_path = self.storage.path("uploads", job.textbook_file)
                source_mime = job.textbook_mime
            async with self.planner_sem, self.storage.read(source_path) as f:
                figures = await docparser.read_figures(self.storage, f, self.storage.path("out", stem, "figures"))
                async for event in agent.plan(source_path, source_mime, f, figures, job.slides, metadata, job.instruction):
                    if isinstance(event, agent.ShikshaCheckpointEvent):
                        await self._process_checkpoint(job, event, "plan", None, out_path)
                    else:
                        await self.jobs.log(job.id, "event", json.loads(event.model_dump_json()))
            await self.jobs.update(job.id, {"status": "creating_slides"})
        elif job.status == "creating_slides":
            await self.jobs.update(job.id, {"message": "Executing presentation outline"})
            stem = pathlib.Path(job.textbook_file).stem
            out_path = self.storage.path("out", stem, f"{job.id}.pptx")
            if not await self.storage.exists(out_path):
                prs = template.Templates.new_presentation()
                await docparser.save_pptx(self.storage, prs, out_path)
            else:
                async with self.storage.read(out_path) as f:
                    prs = await docparser.load_pptx(await f.read())

            outline = agent.PresentationOutline.model_validate(job.metadata["plan"]["outline"])
            metadata = job.metadata.get("design", {})
            if not isinstance(metadata, dict):
                metadata = {}

            if transform_path := job.metadata.get("transform_path", None):
                source_path = self.storage.path("out", stem, transform_path)
            else:
                source_path = self.storage.path("uploads", job.textbook_file)

            figures_dir = self.storage.path("out", stem, "figures")
            async with self.designer_sem, self.storage.read(source_path) as f:
                async for event in agent.design(self.storage, prs, f, figures_dir, outline, metadata, job.instruction):
                    if isinstance(event, agent.ShikshaCheckpointEvent):
                        await self._process_checkpoint(job, event, "design", prs, out_path)
                    else:
                        await self.jobs.log(job.id, "event", json.loads(event.model_dump_json()))

            await docparser.save_pptx(self.storage, prs, out_path)
            failed_slides = metadata.get("failed_slides", [])
            slides_created = metadata.get("slides_created", 0)
            if len(failed_slides) > 0:
                await self.jobs.update(job.id, {"status": "adding_media", "message": f"Slides created ({slides_created} total, {len(failed_slides)} with warnings). Enriching content"})
            else:
                await self.jobs.update(job.id, {"status": "adding_media", "message": f"All {slides_created} slides created successfully. Enriching content"})
        elif job.status == "adding_media":
            await self.jobs.update(job.id, {"message": "Adding media"})
            stem = pathlib.Path(job.textbook_file).stem
            out_path = self.storage.path("out", stem, f"{job.id}.pptx")
            async with self.finalizer_sem:
                async with self.storage.read(out_path) as f:
                    prs = await docparser.load_pptx(await f.read())

                metadata = job.metadata.get("finalize", {})
                if not isinstance(metadata, dict):
                    metadata = {}

                async for event in agent.finalize(prs, metadata):
                    if isinstance(event, agent.ShikshaCheckpointEvent):
                        await self._process_checkpoint(job, event, "finalize", prs, out_path)
                    else:
                        await self.jobs.log(job.id, "event", json.loads(event.model_dump_json()))
                await docparser.save_pptx(self.storage, prs, out_path)
            await self.jobs.update(job.id, {"status": "quality_check"})
        elif job.status == "quality_check":
            stem = pathlib.Path(job.textbook_file).stem
            out_path = self.storage.path("out", stem, f"{job.id}.pptx")
            async with self.storage.read(out_path) as f:
                prs = await docparser.load_pptx(await f.read())

                quality_score = 1.0
                quality_issues = []

                # check 1: minimum slide count
                if len(prs.slides) < 3:
                    quality_score -= 0.3
                    quality_issues.append("Presentation has very few slides")

                # check 2: slide creation success
                design_metadata = job.metadata.get("design", {})
                if isinstance(design_metadata, dict):
                    slides_created = design_metadata.get("slides_created", 0)
                    failed_slides = design_metadata.get("failed_slides", [])

                    if slides_created > 0:
                        quality_score = min(1.0, quality_score + 0.1)

                    if len(failed_slides) > 0:
                        quality_score -= 0.1 * len(failed_slides)
                        quality_issues.append(f"{len(failed_slides)} slide(s) had creation issues")

                # check 3: video integration
                finalize_metadata = job.metadata.get("finalize", {})
                if isinstance(finalize_metadata, dict):
                    videos_added = len(finalize_metadata.get("relevant_videos", []))
                    if videos_added > 0:
                        quality_score = min(1.0, quality_score + 0.1)
                    else:
                        quality_issues.append("No educational videos added")

                # Check 4: Content enrichment
                if isinstance(design_metadata, dict):
                    retry_count = design_metadata.get("retry_count", 0)
                    if retry_count > 2:
                        quality_score -= 0.05
                        quality_issues.append("Multiple retries required during creation")

                # Ensure score stays in valid range
                quality_score = max(0.0, min(1.0, quality_score))

                await self.jobs.update(job.id, {
                    "status": "complete",
                    "message": f"Presentation complete! Quality score: {int(quality_score * 100)}%",
                    "metadata.quality": {
                        "score": quality_score,
                        "issues": quality_issues,
                        "total_slides": len(prs.slides),
                        "completion_time": datetime.now().isoformat()
                    }
                })
        elif job.status == "error":
            if "error" in job.metadata:
                if job.metadata["error"]["attempt"] > self.max_auto_retries:
                    self.logger.warning(f"No longer attempting automatic recovery for job {job.id} ({job.metadata['error']['attempt']} attempts)")
                    if job.metadata["error"]["attempting_recovery"]:
                        await self.jobs.update(job.id, {"metadata.error.attempting_recovery": False})
                    return

                if not job.metadata["error"]["attempting_recovery"]:
                    await self.jobs.update(job.id, {"metadata.error.attempting_recovery": True})

                next_attempt = job.metadata["error"]["next_attempt"]
                delay = max(1, next_attempt - time.time())
                self.logger.info(f"Attempting recovery for job {job.id} after {delay} seconds")
                await asyncio.sleep(delay)

            await self.jobs.update(job.id, {"status": "init", "message": "Retrying after error - attempting recovery"})


def new_default():
    storage = Storage(settings.pres_storage_filesystem, settings.pres_storage_root, settings.pres_storage_options)
    jobs = JobManager(settings.pres_mongodb_url)
    return PresentationService(
        storage,
        jobs,
        settings.pres_do_transform,
        settings.pres_max_auto_retries,
        settings.pres_max_tasks_planner,
        settings.pres_max_tasks_designer,
        settings.pres_max_tasks_finalizer,
    )
