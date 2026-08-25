from dataclasses import dataclass, field
import html
import inspect
import io
import re
from PIL import Image, ImageColor
from pptx import presentation, Presentation
from pptx.slide import Slide
from pptx.util import Emu, Inches, Length, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.text.text import TextFrame
from typing import Annotated, Any, Literal
from pydantic import BaseModel, Field, HttpUrl

HexColor = Annotated[str, Field(pattern=r"^#[0-9A-Fa-f]{6}$", examples=["#FF0000", "#0078D4"])]

SLIDE_WIDTH = Emu(12192000)
SLIDE_HEIGHT = Emu(6858000)
SLIDE_WIDTH_INCHES = SLIDE_WIDTH.inches
SLIDE_HEIGHT_INCHES = SLIDE_HEIGHT.inches
SLIDE_MARGIN_X = SLIDE_WIDTH_INCHES * 0.05
SLIDE_CONTENT_WIDTH = SLIDE_WIDTH_INCHES - (SLIDE_MARGIN_X * 2)
SLIDE_CENTER_CONTENT_LEFT = SLIDE_WIDTH_INCHES * 0.1
SLIDE_CENTER_CONTENT_WIDTH = SLIDE_WIDTH_INCHES - (SLIDE_CENTER_CONTENT_LEFT * 2)

class TextFormatting(BaseModel):
    bold: bool = Field(default=False, description="Set text to bold.")
    italic: bool = Field(default=False, description="Set text to italic.")
    underline: bool = Field(default=False, description="Underline the text.")
    font_color: HexColor | None = Field(default=None, description="Font color in hex format.")
    font_size: int | None = Field(default=None, description="Font size in points.", ge=8, le=96)
    alignment: Literal["left", "center", "right", "justify"] | None = Field(default=None, description="Paragraph alignment.")
    bullet_level: int = Field(default=0, description="Bullet level indentation (0 for top level).", ge=0, le=5)

class TextBlock(BaseModel):
    value: str = Field(description="A single text paragraph or sentence to be displayed as plain body text.", max_length=1000, min_length=1)
    formatting: TextFormatting | None = Field(default=None, description="Formatting options for the text.")

class ListBlock(BaseModel):
    items: list[str] = Field(description="A list of bullet or numbered items to be rendered in sequence.")
    numbered: bool = Field(default=False, description="If True, the list will be rendered as an ordered (numbered) list; if False, it will be an unordered (bullet) list.")

ContentType = str | TextBlock | ListBlock | list[TextBlock | ListBlock]

class WelcomeSlideData(BaseModel):
    title: str = Field(description="Main title text for the slide.", max_length=60, min_length=1, examples=["Welcome to Our Presentation", "Introduction to AI"])
    subtitle: str = Field(description="Subtitle text for the slide.", max_length=120, min_length=1, examples=["A comprehensive overview", "Presented by John Doe"])
    title_font_size: int = Field(default=44, description="Font size for the title in points.", ge=24, le=96)
    subtitle_font_size: int = Field(default=20, description="Font size for the subtitle in points.", ge=12, le=48)
    title_bold: bool = Field(default=True, description="Whether the title should be bold.")
    title_color: HexColor | None = Field(default=None, description="Title color in hex format")
    subtitle_color: HexColor | None = Field(default=None, description="Subtitle color in hex format")
    background_color: HexColor | None = Field(default=None, description="Slide background color in hex format")
    accent_color: HexColor | None = Field(default=None, description="Accent color in hex format")


class BodySlideData(BaseModel):
    heading: str = Field(description="Heading text for the body slide.", max_length=70, min_length=1, examples=["Main Topic", "Key Concepts", "Overview"])
    content: ContentType = Field(description="Main content for the slide. Accepts nested types.", examples=[
        # Example 1: Plain string
        "Social structure: Turks and locals, Bhakti & Sufi influence",

        # Example 2: TextBlock object
        TextBlock(value="Economic policies: improved agriculture, trade routes",
                  formatting=TextFormatting(bold=True, font_color="#1E90FF", font_size=20, alignment="left", bullet_level=0)),

        # Example 3: ListBlock object
        ListBlock(items=[
            "Markets regulated under Khilji",
            "Growth of crafts, trade with foreign merchants",
            "Expansion of trade routes"
        ], numbered=False),

        # Example 4: Nested structure
        [
            "Bhakti & Sufi influence",
            TextBlock(value="Economic policies strengthened agriculture",
                      formatting=TextFormatting(bold=True, font_color="#FF4500")),
            ListBlock(items=[
                "Markets regulated",
                "Growth of crafts",
                "Trade with foreign merchants"
            ], numbered=True)
        ]
    ])
    heading_color: HexColor | None = Field(default=None, description="Heading color in hex format")
    background_color: HexColor | None = Field(default=None, description="Slide background color in hex format")
    accent_color: HexColor | None = Field(default=None, description="Accent color for decorative elements in hex format")


class LayoutBox(BaseModel):
    left: float = Field(description="Left pos of box in inches.", ge=0, le=SLIDE_WIDTH_INCHES)
    top: float = Field(description="Top pos of box in inches.", ge=0, le=SLIDE_HEIGHT_INCHES)
    width: float = Field(description="Width of box in inches.", ge=0.5, le=SLIDE_WIDTH_INCHES)
    height: float = Field(description="Height of box in inches.", ge=0.5, le=SLIDE_HEIGHT_INCHES)


class BodySlideWithImageData(BodySlideData):
    image_path: Any = Field(description="Path or URL to the image to insert.")
    image_placeholder_text: str = Field(default="[Image Placeholder]", description="Text to display if no image is provided or fails to load.", max_length=100, examples=["[Image Placeholder]", "[Diagram]", "[Photo]"])
    image_link: HttpUrl | None = Field(default=None, description="Optional hyperlink for the image.")
    image_caption: str | None = Field(default=None, description="Optional caption text below the image", max_length=200, examples=["Figure 1: Product overview", "Team photo 2024"])
    title_font_size: int = Field(default=32, description="Font size for the title text in points.", ge=12, le=72)
    content_pos: LayoutBox = Field(default=LayoutBox(left=SLIDE_MARGIN_X, top=1.5, width=6.0, height=5.0), description="Position of content box")
    image_pos: LayoutBox = Field(default=LayoutBox(left=7.33, top=1.5, width=5.33, height=5.0), description="Position of image box")
    add_border: bool = Field(default=False, description="Add a border around the image")
    border_color: HexColor | None = Field(default="#0078D4", description="Border color in hex format")
    shadow: bool = Field(default=False, description="Add shadow effect to image")
    image_fit: Literal["contain", "cover", "stretch"] = Field(default="contain", description="How the image should fit its box")


class BodySlideWithVideoData(BodySlideData):
    video: HttpUrl = Field(description="URL of the video to open when the image is clicked.")
    thumbnail_image: Any = Field(description="URL of the thumbnail image to use for the video.")
    thumbnail_width: float = Field(default=5.7, description="Width of video thumbnail box in inches.", ge=1.0, le=10.0)
    thumbnail_height: float = Field(default=3.21, description="Height of video thumbnail box in inches (16:9 aspect ratio recommended).", ge=0.5, le=6.0)
    maintain_aspect_ratio: bool = Field(default=True, description="Maintain original aspect ratio of thumbnail")


class ThankYouSlideData(BaseModel):
    main_text: str = Field(default="Thank you!", description="Main 'Thank You' message text.", max_length=100, examples=["Thank you!", "Thanks for your attention!", "Questions?"])
    contact_info: str | None = Field(default=None, description="Optional contact information text.", max_length=300, examples=["contact@shikshacopilot.in | +91-543-210-1234", "Visit us at shikshacopilot.in"])
    main_text_font_size: int = Field(default=54, description="Font size for main text in points.", ge=24, le=96)
    contact_text_font_size: int = Field(default=18, description="Font size for contact information in points.", ge=10, le=36)
    main_text_pos: LayoutBox = Field(default=LayoutBox(left=SLIDE_CENTER_CONTENT_LEFT, top=2.5, width=SLIDE_CENTER_CONTENT_WIDTH, height=2.0), description="Position of main text box")
    contact_pos: LayoutBox = Field(default=LayoutBox(left=SLIDE_CENTER_CONTENT_LEFT, top=5.0, width=SLIDE_CENTER_CONTENT_WIDTH, height=1.5), description="Position of contact text box")
    background_color: HexColor = Field(default="#F8FAFC", description="Slide background color")
    main_text_color: HexColor = Field(default="#0F172A", description="Main text color")
    contact_text_color: HexColor = Field(default="#475569", description="Contact text color")
    accent_color: HexColor = Field(default="#2563EB", description="Accent color")


class FullImageSlideData(BaseModel):
    """Full-bleed image slide with text overlay - image fills entire slide."""
    image_path: Any = Field(description="Path or URL to the background image.")
    title: str = Field(description="Main title text overlaid on image.", max_length=50, min_length=1, examples=["Our Vision", "The Future"])
    subtitle: str | None = Field(default=None, description="Optional subtitle text.", max_length=120, examples=["Building tomorrow, today", "Innovation at scale"])
    title_position: Literal["top", "center", "bottom"] = Field(default="center", description="Title position")
    text_color: HexColor = Field(default="#FFFFFF", description="Text color in hex format.")
    add_overlay: bool = Field(default=True, description="Add semi-transparent overlay for text readability")
    overlay_color: HexColor = Field(default="#000000", description="Overlay color in hex format")
    overlay_opacity: float = Field(default=0.4, description="Overlay opacity (0.0-1.0)", ge=0.0, le=1.0)
    title_font_size: int = Field(default=54, description="Font size for title in points.", ge=24, le=96)
    subtitle_font_size: int = Field(default=24, description="Font size for subtitle in points.", ge=12, le=48)


class IconItem(BaseModel):
    """Icon with text description."""
    icon_text: str = Field(description="Text to display as icon (emoji or single character)", max_length=5, min_length=1, examples=["🚀", "💡", "⭐", "✓"])
    title: str = Field(description="Title for this icon item", max_length=30, min_length=1, examples=["Innovation", "Quality", "Speed"])
    description: str = Field(description="Description text", max_length=120, min_length=1, examples=["We innovate constantly", "Quality is our priority"])


class IconSlideData(BaseModel):
    """Icon-based slide with icons and short descriptions."""
    heading: str = Field(description="Main heading for the slide.", max_length=70, min_length=1, examples=["Our Values", "Key Features", "Core Principles"])
    icons: list[IconItem] = Field(description="List of icon items (2-4 recommended)")
    layout: Literal["grid", "horizontal", "vertical"] = Field(default="grid", description="Layout style")
    icon_color: HexColor = Field(default="#0078D4", description="Color for icons in hex format")
    background_color: HexColor | None = Field(default=None, description="Slide background color")
    heading_color: HexColor = Field(default="#0F172A", description="Heading color")
    accent_color: HexColor = Field(default="#F57C00", description="Accent color")


class HighlightBox(BaseModel):
    """Colored box with content."""
    title: str = Field(description="Box title", max_length=25, min_length=1, examples=["Key Feature", "Important Point"])
    content: str = Field(description="Box content", max_length=200, min_length=1, examples=["This is the main content of the highlight box"])
    color: HexColor = Field(default="#0078D4", description="Box color in hex format")


class HighlightBoxSlideData(BaseModel):
    """Slide with colored highlight boxes for key information."""
    heading: str = Field(description="Main heading for the slide.", max_length=70, min_length=1, examples=["Key Features", "Main Benefits"])
    boxes: list[HighlightBox] = Field(description="List of highlight boxes (2-4 recommended)")
    layout: Literal["grid", "horizontal", "vertical"] = Field(default="grid", description="Layout style")
    background_color: HexColor = Field(default="#F8FAFC", description="Slide background color")
    heading_color: HexColor = Field(default="#0F172A", description="Heading color")
    accent_color: HexColor = Field(default="#2563EB", description="Fallback accent color")
    title_font_size: int = Field(default=34, description="Font size for heading in points.", ge=18, le=60)


class BigNumberSlideData(BaseModel):
    """Slide with large number/statistic and context."""
    number: str = Field(description="The big number or statistic", max_length=12, min_length=1, examples=["95%", "1,000+", "$5M", "2.5x"])
    context: str = Field(description="Context or explanation for the number", max_length=100, min_length=1, examples=["Customer satisfaction rate", "Active users worldwide"])
    subtitle: str | None = Field(default=None, description="Optional additional context", max_length=120, examples=["Based on 2024 survey results"])
    number_color: HexColor = Field(default="#0078D4", description="Color for the number in hex format")
    background_color: HexColor | None = Field(default=None, description="Slide background color")
    context_color: HexColor = Field(default="#1F2937", description="Context text color")
    accent_color: HexColor = Field(default="#F57C00", description="Accent color")


class ProcessStep(BaseModel):
    """Single step in a process."""
    step_number: int = Field(description="Step number", ge=1, le=20)
    title: str = Field(description="Step title", max_length=23, min_length=1, examples=["Research", "Design", "Implementation"])
    description: str = Field(description="Step description", max_length=150, min_length=1, examples=["Conduct market research and gather requirements"])


class ProcessFlowSlideData(BaseModel):
    """Process flow slide showing sequential steps."""
    heading: str = Field(description="Main heading for the process.", max_length=70, min_length=1, examples=["Our Development Process", "Project Timeline"])
    steps: list[ProcessStep] = Field(description="List of process steps (3-5 recommended)")
    flow_color: HexColor = Field(default="#0078D4", description="Color for flow elements in hex format")
    background_color: HexColor | None = Field(default=None, description="Slide background color")
    heading_color: HexColor = Field(default="#0F172A", description="Heading color")
    accent_color: HexColor = Field(default="#F57C00", description="Secondary accent color")
    title_font_size: int = Field(default=34, description="Font size for heading in points.", ge=18, le=60)


class TwoColumnSlideData(BaseModel):
    """Two-column layout with heading and separate content areas."""
    heading: str = Field(description="Heading text for the slide.", max_length=70, min_length=1, examples=["Comparison", "Before & After"])
    left_content: ContentType = Field(description="Content for the left column.")
    right_content: ContentType = Field(description="Content for the right column.")
    title_font_size: int = Field(default=32, description="Font size for the title in points.", ge=12, le=72)
    column_gap: float = Field(default=0.4, description="Gap between columns in inches.", ge=0.1, le=2.0)
    background_color: HexColor = Field(default="#F9FBFF", description="Slide background color in hex format")
    heading_color: HexColor = Field(default="#0F172A", description="Heading color in hex format")
    accent_color: HexColor = Field(default="#0078D4", description="Accent color for decorative elements in hex format")
    left_background: HexColor | None = Field(default=None, description="Background color for left column")
    right_background: HexColor | None = Field(default=None, description="Background color for right column")
    add_divider: bool = Field(default=False, description="Add vertical divider between columns")


class QuoteSlideData(BaseModel):
    """Quote slide with large centered quote and attribution."""
    quote_text: str = Field(description="The main quote text.", max_length=250, min_length=1, examples=["Innovation distinguishes between a leader and a follower."])
    attribution: str | None = Field(default=None, description="Author or source of the quote.", max_length=50, examples=["Steve Jobs", "Albert Einstein"])
    quote_font_size: int = Field(default=36, description="Font size for the quote in points.", ge=18, le=72)
    attribution_font_size: int = Field(default=20, description="Font size for attribution in points.", ge=12, le=48)
    background_color: HexColor = Field(default="#F8FAFC", description="Slide background color")
    quote_color: HexColor = Field(default="#0F172A", description="Quote text color")
    attribution_color: HexColor = Field(default="#475569", description="Attribution text color")
    accent_color: HexColor = Field(default="#2563EB", description="Accent color")


class ComparisonSlideData(BaseModel):
    """Side-by-side comparison slide with two titled sections."""
    heading: str = Field(description="Main heading for the comparison.", max_length=70, min_length=1, examples=["Traditional vs Modern", "Before & After"])
    left_title: str = Field(description="Title for the left comparison section.", max_length=50, min_length=1, examples=["Traditional Approach", "Before"])
    left_content: ContentType = Field(description="Content for the left section.")
    right_title: str = Field(description="Title for the right comparison section.", max_length=50, min_length=1, examples=["Modern Approach", "After"])
    right_content: ContentType = Field(description="Content for the right section.")
    title_font_size: int = Field(default=32, description="Font size for main heading in points.", ge=12, le=72)
    section_title_font_size: int = Field(default=24, description="Font size for section titles in points.", ge=12, le=48)
    background_color: HexColor = Field(default="#F8FAFC", description="Slide background color")
    heading_color: HexColor = Field(default="#0F172A", description="Heading color")
    left_color: HexColor = Field(default="#2563EB", description="Left section accent color")
    right_color: HexColor = Field(default="#F57C00", description="Right section accent color")


class TimelineEvent(BaseModel):
    date: str = Field(description="Date or time of the event.", max_length=30, min_length=1, examples=["2024", "Q1 2024", "January 2024"])
    description: str = Field(description="Description of the event.", max_length=150, min_length=1, examples=["Product launch", "Major milestone achieved"])


class TimelineSlideData(BaseModel):
    """Timeline slide with chronological events."""
    heading: str = Field(description="Heading text for the timeline.", max_length=70, min_length=1, examples=["Company History", "Project Timeline"])
    events: list[TimelineEvent] = Field(description="List of timeline events.")
    title_font_size: int = Field(default=32, description="Font size for the title in points.", ge=12, le=72)
    event_font_size: int = Field(default=18, description="Font size for events in points.", ge=10, le=36)
    background_color: HexColor = Field(default="#F8FAFC", description="Slide background color")
    heading_color: HexColor = Field(default="#0F172A", description="Heading color")
    accent_color: HexColor = Field(default="#2563EB", description="Timeline accent color")


class ImageGridItem(BaseModel):
    image: Any = Field(description="Path, URL, or binary stream for the image.")
    caption: str | None = Field(default=None, description="Optional caption for the image.", max_length=100, examples=["Product showcase", "Team photo"])


class ImageGridSlideData(BaseModel):
    """Grid layout with multiple images and optional captions."""
    heading: str = Field(description="Heading text for the slide.", max_length=70, min_length=1, examples=["Our Gallery", "Product Showcase"])
    images: list[ImageGridItem] = Field(description="List of image items with optional captions.")
    grid_columns: int = Field(default=2, description="Number of columns in the grid.", ge=1, le=4)
    title_font_size: int = Field(default=32, description="Font size for the title in points.", ge=12, le=72)
    background_color: HexColor = Field(default="#F8FAFC", description="Slide background color")
    heading_color: HexColor = Field(default="#0F172A", description="Heading color")
    accent_color: HexColor = Field(default="#2563EB", description="Accent color")
    image_fit: Literal["contain", "cover", "stretch"] = Field(default="contain", description="How images should fit each grid cell")


class SectionDividerSlideData(BaseModel):
    """Section divider slide with large centered text."""
    section_title: str = Field(description="Main section title text.", max_length=39, min_length=1, examples=["Introduction", "Chapter 1", "Part Two"])
    section_number: str | None = Field(default=None, description="Optional section number or label.", max_length=10, examples=["01", "Chapter 1", "Part I"])
    background_color: HexColor = Field(default="#4472C4", description="Background color in hex format.")
    text_color: HexColor = Field(default="#FFFFFF", description="Text color in hex format.")
    title_font_size: int = Field(default=54, description="Font size for section title in points.", ge=24, le=96)


class KeyTakeawaysSlideData(BaseModel):
    """Key takeaways slide with numbered or bulleted highlights."""
    heading: str = Field(default="Key Takeaways", description="Heading text for the slide.", max_length=70, examples=["Key Takeaways", "Summary", "Main Points"])
    takeaways: list[str] = Field(description="List of key takeaway points.", min_length=1, max_length=6)
    numbered: bool = Field(default=True, description="Whether to use numbered list.")
    title_font_size: int = Field(default=36, description="Font size for the title in points.", ge=12, le=72)
    takeaway_font_size: int = Field(default=22, description="Font size for takeaways in points.", ge=10, le=48)
    layout: Literal["auto", "list", "grid"] = Field(default="auto", description="Takeaway layout")
    background_color: HexColor = Field(default="#F8FAFC", description="Slide background color")
    heading_color: HexColor = Field(default="#0F172A", description="Heading color")
    accent_color: HexColor = Field(default="#2563EB", description="Accent color")


class AgendaItem(BaseModel):
    topic: str = Field(description="Agenda topic or activity.", max_length=100, min_length=1, examples=["Introduction and Welcome", "Q&A Session"])
    time: str | None = Field(default=None, description="Optional time allocation for the agenda item.", max_length=30, examples=["10:00 AM", "30 minutes", "15 min"])


class AgendaSlideData(BaseModel):
    """Agenda slide with list of topics and optional time allocations."""
    heading: str = Field(default="Agenda", description="Heading text for the slide.", max_length=70, examples=["Agenda", "Today's Schedule", "Meeting Outline"])
    items: list[AgendaItem] = Field(description="List of agenda items with optional time allocations.")
    title_font_size: int = Field(default=36, description="Font size for the title in points.", ge=12, le=72)
    item_font_size: int = Field(default=20, description="Font size for agenda items in points.", ge=10, le=36)
    background_color: HexColor = Field(default="#F8FAFC", description="Slide background color")
    heading_color: HexColor = Field(default="#0F172A", description="Heading color")
    accent_color: HexColor = Field(default="#2563EB", description="Accent color")


class QuestionSlideData(BaseModel):
    """Interactive question slide to engage students and promote critical thinking."""
    question: str = Field(description="The main question to ask students", max_length=200, min_length=1, examples=["What factors led to the Industrial Revolution?", "How does photosynthesis work?"])
    question_type: Literal["open", "multiple_choice", "true_false", "poll"] = Field(default="open", description="Type of question")
    options: list[str] | None = Field(default=None, description="Answer options for multiple choice or poll questions")
    hint: str | None = Field(default=None, description="Optional hint or guidance for students", max_length=120, examples=["Think about the time period", "Consider the process steps"])
    background_color: HexColor | None = Field(default="#F0F8FF", description="Slide background color")
    question_color: HexColor = Field(default="#0078D4", description="Question text color")
    icon: str = Field(default="❓", description="Emoji icon for the question", max_length=5)


class DiscussionSlideData(BaseModel):
    """Discussion prompt slide with thought-provoking questions for classroom engagement."""
    heading: str = Field(description="Main heading for the discussion", max_length=70, min_length=1, examples=["Let's Discuss", "Group Discussion", "Think About It"])
    discussion_prompts: list[str] = Field(description="List of discussion questions or prompts (2-4 recommended)")
    discussion_type: Literal["group", "pair", "class", "think_pair_share"] = Field(default="group")
    time_allocation: str | None = Field(default=None, description="Suggested time for discussion", max_length=30, examples=["10 minutes", "15 min", "5-10 minutes"])
    background_color: HexColor | None = Field(default="#FFF8DC", description="Slide background color")
    heading_color: HexColor = Field(default="#D32F2F", description="Heading color")
    accent_color: HexColor = Field(default="#F57C00", description="Accent color")
    icon: str = Field(default="💬", description="Emoji icon for discussion", max_length=5)


class ActivitySlideData(BaseModel):
    """Activity slide with instructions for hands-on student engagement."""
    activity_title: str = Field(description="Title of the activity", max_length=70, min_length=1, examples=["Hands-On Experiment", "Group Activity", "Practice Exercise"])
    instructions: list[str] = Field(description="Step-by-step instructions (3-5 steps recommended)")
    materials_needed: list[str] | None = Field(default=None, description="Materials or resources needed")
    time_required: str | None = Field(default=None, description="Estimated time", max_length=30, examples=["15 minutes", "20-30 min", "1 hour"])
    group_size: str | None = Field(default=None, description="Suggested group size", max_length=30, examples=["Pairs", "Groups of 4", "Individual"])
    background_color: HexColor | None = Field(default="#F0FFF0", description="Slide background color")
    title_color: HexColor = Field(default="#2E7D32", description="Title color")
    accent_color: HexColor = Field(default="#F57C00", description="Accent color")
    icon: str = Field(default="✏️", description="Emoji icon for activity", max_length=5)


class RealWorldConnectionSlideData(BaseModel):
    """Slide connecting lesson content to real-world applications and relevance."""
    heading: str = Field(description="Main heading", max_length=70, min_length=1, examples=["Real-World Connection", "Why This Matters", "Practical Applications"])
    connection_title: str = Field(description="Title of the real-world connection", max_length=70, min_length=1, examples=["How This Applies to Your Life", "Industry Applications"])
    description: str = Field(description="Description of how the content applies to real life", max_length=300, min_length=1, examples=["This concept is used daily in modern technology..."])
    examples: list[str] = Field(description="Real-world examples (2-4 recommended)")
    call_to_action: str | None = Field(default=None, description="Optional call to action or reflection prompt", max_length=120, examples=["Think about how you use this in your daily life"])
    background_color: HexColor | None = Field(default="#FFF5E6", description="Slide background color")
    heading_color: HexColor = Field(default="#F57C00", description="Heading color")
    accent_color: HexColor = Field(default="#2563EB", description="Accent color")
    icon: str = Field(default="🌍", description="Emoji icon for real-world connection", max_length=5)


class StorySlideData(BaseModel):
    """Narrative/story slide to make content memorable through storytelling."""
    story_title: str = Field(description="Title of the story or narrative", max_length=70, min_length=1, examples=["The Story of Discovery", "A Historical Perspective"])
    story_content: str = Field(description="The story or narrative text (keep engaging and concise)", max_length=600, min_length=1, examples=["Once upon a time, in a small laboratory..."])
    story_type: Literal["historical", "personal", "case_study", "anecdote"] = Field(default="historical")
    moral_or_lesson: str | None = Field(default=None, description="Key takeaway or lesson from the story", max_length=150, examples=["Persistence leads to breakthrough discoveries"])
    background_color: HexColor | None = Field(default="#FFF0F5", description="Slide background color")
    title_color: HexColor = Field(default="#8B5CF6", description="Title color")
    accent_color: HexColor = Field(default="#F57C00", description="Accent color")
    icon: str = Field(default="📖", description="Emoji icon for story", max_length=5)


SlideCategory = Literal["content", "engagement", "structural", "visual"]


@dataclass(frozen=True)
class SlideMeta:
    categories: set[SlideCategory]
    tags: set[Literal["supports_images"]] = field(default_factory=set)


SLIDE_META: dict[str, SlideMeta] = {}
def slide(meta: SlideMeta):
    def decorator(fn):
        SLIDE_META[fn.__name__] = meta
        return fn
    return decorator


class Templates:

    _BASE_SLIDE_WIDTH = SLIDE_WIDTH
    _BASE_SLIDE_HEIGHT = SLIDE_HEIGHT
    _BASE_SLIDE_WIDTH_INCHES = SLIDE_WIDTH_INCHES
    _BASE_SLIDE_HEIGHT_INCHES = SLIDE_HEIGHT_INCHES
    _MARGIN_X = SLIDE_MARGIN_X
    _CONTENT_WIDTH = SLIDE_CONTENT_WIDTH
    _CENTER_CONTENT_LEFT = SLIDE_CENTER_CONTENT_LEFT
    _CENTER_CONTENT_WIDTH = SLIDE_CENTER_CONTENT_WIDTH

    @staticmethod
    def new_presentation(slide_width: Length = _BASE_SLIDE_WIDTH, slide_height: Length = _BASE_SLIDE_HEIGHT) -> presentation.Presentation:
        prs = Presentation()
        prs.slide_width = slide_width
        prs.slide_height = slide_height
        return prs

    def __init__(self, prs: presentation.Presentation | None = None):
        self.prs: presentation.Presentation = prs or Templates.new_presentation()

    @slide(SlideMeta(categories={"structural", "engagement"}))
    def welcome(self, data: WelcomeSlideData):
        """
        Create the opening title slide for the full deck.

        Choose this only at the start of the presentation, with a short hook title
        and a scene-setting subtitle. For chapter or topic transitions, choose the
        section divider tool instead.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string((data.background_color or "#F8FAFC").replace("#", ""))
        accent_rgb = RGBColor.from_string((data.accent_color or data.title_color or "#2563EB").replace("#", ""))

        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._MARGIN_X), Inches(1.05), Inches(1.7), Inches(0.07))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_rgb
        accent_bar.line.fill.background()

        title_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X), Inches(2.05), Inches(self._CONTENT_WIDTH), Inches(1.55)).text_frame
        title_frame.text = data.title
        title_frame.word_wrap = True
        title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(min(data.title_font_size, 48 if len(data.title) > 48 else data.title_font_size))
        title_para.font.bold = data.title_bold
        title_para.alignment = PP_ALIGN.LEFT

        if data.title_color:
            title_rgb = RGBColor.from_string(data.title_color.replace('#', ''))
            title_para.font.color.rgb = title_rgb
        else:
            title_para.font.color.rgb = RGBColor(15, 23, 42)

        subtitle_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X), Inches(4.25), Inches(self._CONTENT_WIDTH * 0.72), Inches(1.0)).text_frame
        subtitle_frame.text = data.subtitle
        subtitle_frame.word_wrap = True
        subtitle_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(data.subtitle_font_size)
        subtitle_para.alignment = PP_ALIGN.LEFT
        subtitle_para.line_spacing = 1.12

        if data.subtitle_color:
            subtitle_rgb = RGBColor.from_string(data.subtitle_color.replace('#', ''))
            subtitle_para.font.color.rgb = subtitle_rgb
        else:
            subtitle_para.font.color.rgb = RGBColor(71, 85, 105)

        for shape in slide.shapes:
            if shape.has_text_frame:
                self._fit_text(shape.text_frame)

        return slide


    @slide(SlideMeta(categories={"content"}))
    def body(self, data: BodySlideData):
        """
        Create a plain text-focused slide with a heading and content area.

        Choose this only when no stronger visual, sequence, comparison, metric, or
        activity format fits. Keep it to one main idea and 3-4 concise points.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # Apply background color if specified, otherwise use subtle gradient
        if data.background_color:
            background = slide.background
            fill = background.fill
            fill.solid()
            rgb = RGBColor.from_string(data.background_color.replace('#', ''))
            fill.fore_color.rgb = rgb
        else:
            # Add very subtle gradient for depth
            background = slide.background
            fill = background.fill
            fill.gradient()
            fill.gradient_angle = 90.0
            fill.gradient_stops[0].color.rgb = RGBColor(250, 250, 252)
            fill.gradient_stops[1].color.rgb = RGBColor(255, 255, 255)

        title = slide.shapes.add_textbox(Inches(self._MARGIN_X), Inches(0.5), Inches(self._CONTENT_WIDTH), Inches(0.8))
        title.text = data.heading

        # Enhanced title styling
        title_para = title.text_frame.paragraphs[0]
        title_para.font.size = Pt(36)  # Slightly larger
        title_para.font.bold = True
        title_para.font.name = 'Calibri'

        # Apply heading color if specified
        if data.heading_color:
            rgb = RGBColor.from_string(data.heading_color.replace('#', ''))
            title_para.font.color.rgb = rgb
        else:
            title_para.font.color.rgb = RGBColor(0, 51, 102)  # Professional dark blue

        text_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT), Inches(1.55), Inches(self._CENTER_CONTENT_WIDTH), Inches(4.9)).text_frame

        # Add padding to content area
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.1)

        self._render_content(text_frame, data.content)

        # Enhanced accent decorative elements
        accent_color = data.accent_color if data.accent_color else '#0078D4'
        rgb = RGBColor.from_string(accent_color.replace('#', ''))

        # Modern accent bar under title - wider and more prominent
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._MARGIN_X), Inches(1.25), Inches(4.0), Inches(0.08))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = rgb
        accent_bar.line.fill.background()

        # Add subtle decorative corner element
        corner_accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._BASE_SLIDE_WIDTH_INCHES - self._MARGIN_X - 0.2), Inches(0.5), Inches(0.2), Inches(1.0))
        corner_accent.fill.solid()
        corner_accent.fill.fore_color.rgb = rgb
        corner_accent.fill.transparency = 0.7
        corner_accent.line.fill.background()

        # Add small decorative circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(self._MARGIN_X * 0.5), Inches(6.8), Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = rgb
        circle.fill.transparency = 0.6
        circle.line.fill.background()

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"content", "visual"}, tags={"supports_images"}))
    def body_with_image(self, data: BodySlideWithImageData):
        """
        Create an explanatory slide with concise text beside one supporting image.

        Best for people, places, objects, diagrams, maps, or examples where the
        image clarifies the lesson. Include attribution in the caption when using
        searched images, and use a border or shadow for visual polish.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide: Slide = self.prs.slides.add_slide(blank_slide_layout)

        # Apply background color if specified
        if data.background_color:
            background = slide.background
            fill = background.fill
            fill.solid()
            rgb = RGBColor.from_string(data.background_color.replace('#', ''))
            fill.fore_color.rgb = rgb

        # Title
        title_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X), Inches(0.5), Inches(self._CONTENT_WIDTH), Inches(0.8)).text_frame
        title_frame.text = data.heading
        title_frame.paragraphs[0].font.size = Pt(data.title_font_size)
        title_frame.paragraphs[0].font.bold = True

        # Apply heading color if specified
        if data.heading_color:
            rgb = RGBColor.from_string(data.heading_color.replace('#', ''))
            title_frame.paragraphs[0].font.color.rgb = rgb

        # Content area
        content_frame = slide.shapes.add_textbox(Inches(data.content_pos.left), Inches(data.content_pos.top), Inches(data.content_pos.width), Inches(data.content_pos.height)).text_frame
        content_frame.word_wrap = True

        self._render_content(content_frame, data.content)

        image_shape = None
        if data.image_path:
            try:
                image_shape = self._add_picture_fit(slide, data.image_path, Inches(data.image_pos.left), Inches(data.image_pos.top), Inches(data.image_pos.width), Inches(data.image_pos.height), data.image_fit)

                if data.add_border and data.border_color:
                    image_shape.line.color.rgb = RGBColor.from_string(data.border_color.replace('#', ''))
                    image_shape.line.width = Pt(2)

                if data.shadow:
                    shadow = image_shape.shadow
                    shadow.inherit = False
                    shadow.visible = True
                    shadow.style = 1  # Outer shadow
                    shadow.blur_radius = Pt(4)
                    shadow.distance = Pt(3)
                    shadow.angle = 45

            except Exception as e:
                image_shape = self._create_image_placeholder(
                    slide,
                    Inches(data.image_pos.left),
                    Inches(data.image_pos.top),
                    Inches(data.image_pos.width),
                    Inches(data.image_pos.height),
                    f"Error loading image:\\n{str(e)}"
                )
        else:
            image_shape = self._create_image_placeholder(
                slide,
                Inches(data.image_pos.left),
                Inches(data.image_pos.top),
                Inches(data.image_pos.width),
                Inches(data.image_pos.height),
                data.image_placeholder_text
            )

        # Add image caption if provided
        if data.image_caption and image_shape:
            caption_frame = slide.shapes.add_textbox(Inches(data.image_pos.left), Inches(data.image_pos.top + data.image_pos.height + 0.1), Inches(data.image_pos.width), Inches(0.4)).text_frame
            caption_frame.text = data.image_caption
            caption_frame.paragraphs[0].font.size = Pt(12)
            caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            caption_frame.paragraphs[0].font.italic = True
            caption_frame.paragraphs[0].font.color.rgb = RGBColor(96, 94, 92)

        # Hyperlink
        if image_shape is not None and data.image_link:
            image_shape.click_action.hyperlink.address = str(data.image_link)

        # Add accent color decorative element if specified
        if data.accent_color:
            accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._MARGIN_X), Inches(1.2), Inches(2.67), Inches(0.05))
            accent_bar.fill.solid()
            rgb = RGBColor.from_string(data.accent_color.replace('#', ''))
            accent_bar.fill.fore_color.rgb = rgb
            accent_bar.line.fill.background()

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"content", "engagement", "visual"}, tags={"supports_images"}))
    def body_with_video(self, data: BodySlideWithVideoData):
        """Add a body slide with heading, content area, and video placeholder/actual video to the PowerPoint."""
        blank_slide_layout = self.prs.slide_layouts[6]
        slide: Slide = self.prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string((data.background_color or "#F7F9FC").replace("#", ""))
        accent_hex = (data.accent_color or "#2563EB").replace("#", "")
        accent_rgb = RGBColor.from_string(accent_hex)
        heading_rgb = RGBColor.from_string((data.heading_color or "#0F172A").replace("#", ""))

        title_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X), Inches(0.88), Inches(4.75), Inches(1.5)).text_frame
        title_frame.text = data.heading
        title_frame.word_wrap = True
        title_frame.paragraphs[0].font.size = Pt(38 if len(data.heading) < 48 else 32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = heading_rgb
        title_frame.paragraphs[0].line_spacing = 0.94

        content_rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._MARGIN_X), Inches(2.65), Inches(0.85), Inches(0.05))
        content_rule.fill.solid()
        content_rule.fill.fore_color.rgb = accent_rgb
        content_rule.line.fill.background()

        content_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X), Inches(3.1), Inches(4.75), Inches(2.25)).text_frame
        content_frame.margin_left = Inches(0.02)
        content_frame.word_wrap = True
        self._render_content(content_frame, data.content)
        for paragraph in content_frame.paragraphs:
            paragraph.font.size = Pt(17)
            paragraph.font.color.rgb = RGBColor(51, 65, 85)
            paragraph.line_spacing = 1.2

        media_area_left = 6.35
        media_area_top = 1.0
        media_area_width = 5.7
        media_area_height = 4.65
        thumbnail_aspect = self._image_aspect(data.thumbnail_image) or (16.0 / 9.0)
        portrait_like = thumbnail_aspect < 0.95
        letterboxed_wide = thumbnail_aspect >= 1.1 and data.maintain_aspect_ratio

        if portrait_like:
            preview_height = media_area_height
            preview_width = min(media_area_width, preview_height * thumbnail_aspect)
            image_fit = "contain"
        elif letterboxed_wide:
            preview_width = media_area_width
            preview_height = min(media_area_height, preview_width / thumbnail_aspect)
            image_fit = "contain"
        else:
            preview_width = min(media_area_width, 3.55)
            preview_height = media_area_height
            image_fit = "cover"

        preview_left = media_area_left + (media_area_width - preview_width) / 2
        preview_top = media_area_top + (media_area_height - preview_height) / 2

        media_frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(preview_left - 0.05), Inches(preview_top - 0.05), Inches(preview_width + 0.1), Inches(preview_height + 0.1))
        media_frame.fill.solid()
        media_frame.fill.fore_color.rgb = RGBColor(255, 255, 255)
        media_frame.line.fill.background()
        media_frame.shadow.inherit = False
        media_frame.shadow.visible = True
        media_frame.shadow.blur_radius = Pt(8)
        media_frame.shadow.distance = Pt(2)

        try:
            thumbnail = self._add_picture_fit(slide, data.thumbnail_image, Inches(preview_left), Inches(preview_top), Inches(preview_width), Inches(preview_height), image_fit)
        except Exception as e:
            thumbnail = self._create_image_placeholder(slide, Inches(preview_left), Inches(preview_top), Inches(preview_width), Inches(preview_height), f"Error loading video thumbnail:\n{str(e)}")

        play_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(preview_left + preview_width / 2 - 0.36), Inches(preview_top + preview_height / 2 - 0.36), Inches(0.72), Inches(0.72))
        play_bg.fill.solid()
        play_bg.fill.fore_color.rgb = accent_rgb
        play_bg.line.fill.background()

        play = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(preview_left + preview_width / 2 - 0.07), Inches(preview_top + preview_height / 2 - 0.14), Inches(0.22), Inches(0.28))
        play.rotation = 90
        play.fill.solid()
        play.fill.fore_color.rgb = RGBColor(255, 255, 255)
        play.line.fill.background()

        for shape in (media_frame, thumbnail, play_bg, play):
            shape.click_action.hyperlink.address = str(data.video)

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"structural"}))
    def thank_you(self, data: ThankYouSlideData):
        """Add a thank you slide with centered text and optional contact information to the PowerPoint."""
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(data.background_color.replace("#", ""))
        accent_rgb = RGBColor.from_string(data.accent_color.replace("#", ""))

        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._MARGIN_X), Inches(1.45), Inches(1.8), Inches(0.07))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_rgb
        accent_bar.line.fill.background()

        thank_you_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT), Inches(2.45), Inches(self._CENTER_CONTENT_WIDTH), Inches(1.0)).text_frame
        thank_you_frame.text = data.main_text
        thank_you_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        thank_you_frame.paragraphs[0].font.size = Pt(data.main_text_font_size)
        thank_you_frame.paragraphs[0].font.bold = True
        thank_you_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(data.main_text_color.replace("#", ""))

        if data.contact_info:
            contact_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT + 0.65), Inches(4.2), Inches(self._CENTER_CONTENT_WIDTH - 1.3), Inches(1.0)).text_frame
            contact_frame.text = data.contact_info
            contact_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            contact_frame.paragraphs[0].font.size = Pt(data.contact_text_font_size)
            contact_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(data.contact_text_color.replace("#", ""))
            contact_frame.paragraphs[0].line_spacing = 1.12
            contact_frame.word_wrap = True

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"content", "visual"}))
    def two_column(self, data: TwoColumnSlideData):
        """
        Create a slide with two balanced content columns under one heading.

        Best for paired ideas such as cause and effect, before and after notes, or
        parallel categories. For a direct opposition with named sides, choose the
        comparison tool instead.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        background = slide.background
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor.from_string(data.background_color.replace('#', ''))
        heading_rgb = RGBColor.from_string(data.heading_color.replace('#', ''))
        accent_rgb = RGBColor.from_string(data.accent_color.replace('#', ''))

        title_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X), Inches(0.45), Inches(self._CONTENT_WIDTH), Inches(0.9)).text_frame
        title_frame.text = data.heading
        title_frame.paragraphs[0].font.size = Pt(data.title_font_size)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = heading_rgb
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        title_frame.word_wrap = True

        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._MARGIN_X), Inches(1.28), Inches(1.55), Inches(0.05))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_rgb
        accent_bar.line.fill.background()

        column_width = (self._CONTENT_WIDTH - data.column_gap) / 2
        column_top = 1.75
        column_height = 4.85
        left = self._MARGIN_X
        right = self._MARGIN_X + column_width + data.column_gap

        def content_items(entry: ContentType) -> list[str]:
            if isinstance(entry, str):
                return [html.unescape(entry)]
            if isinstance(entry, TextBlock):
                return [html.unescape(entry.value)]
            if isinstance(entry, ListBlock):
                return [html.unescape(item) for item in entry.items]
            return [item for content in entry for item in content_items(content)]

        def render_column(column_left: float, marker: str, marker_rgb: RGBColor, content: ContentType, background_color: str | None):
            if background_color:
                tint = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(column_left - 0.12), Inches(column_top - 0.08), Inches(column_width + 0.24), Inches(column_height + 0.16))
                tint.fill.solid()
                tint.fill.fore_color.rgb = RGBColor.from_string(background_color.replace('#', ''))
                tint.fill.transparency = 0.18
                tint.line.fill.background()

            marker_frame = slide.shapes.add_textbox(Inches(column_left), Inches(column_top), Inches(0.7), Inches(0.45)).text_frame
            marker_frame.text = marker
            marker_frame.paragraphs[0].font.size = Pt(24)
            marker_frame.paragraphs[0].font.bold = True
            marker_frame.paragraphs[0].font.color.rgb = marker_rgb

            rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(column_left + 0.75), Inches(column_top + 0.25), Inches(column_width - 0.75), Inches(0.02))
            rule.fill.solid()
            rule.fill.fore_color.rgb = RGBColor(203, 213, 225)
            rule.line.fill.background()

            items = content_items(content)
            row_gap = 0.2
            row_height = (column_height - 0.8 - (len(items) - 1) * row_gap) / max(len(items), 1)
            row_top = column_top + 0.85
            for item in items:
                dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(column_left + 0.02), Inches(row_top + 0.12), Inches(0.12), Inches(0.12))
                dot.fill.solid()
                dot.fill.fore_color.rgb = marker_rgb
                dot.line.fill.background()

                item_frame = slide.shapes.add_textbox(Inches(column_left + 0.32), Inches(row_top), Inches(column_width - 0.32), Inches(row_height)).text_frame
                item_frame.text = item
                item_frame.word_wrap = True
                item_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                item_frame.margin_left = 0
                item_frame.margin_right = 0
                item_frame.paragraphs[0].font.size = Pt(17 if len(item) < 92 else 15)
                item_frame.paragraphs[0].font.color.rgb = RGBColor(30, 41, 59)
                item_frame.paragraphs[0].line_spacing = 1.15
                row_top += row_height + row_gap

        render_column(left, "01", accent_rgb, data.left_content, data.left_background)
        render_column(right, "02", heading_rgb, data.right_content, data.right_background)

        if data.add_divider:
            divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._MARGIN_X + column_width + (data.column_gap / 2) - 0.01), Inches(column_top + 0.1), Inches(0.02), Inches(column_height - 0.2))
            divider.fill.solid()
            divider.fill.fore_color.rgb = RGBColor(203, 213, 225)
            divider.line.fill.background()

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"content", "engagement"}))
    def quote(self, data: QuoteSlideData):
        """
        Create a reflective slide centered on one quote or source statement.

        Best for a discussion opener, emotional anchor, or pause for interpretation.
        Include attribution when known. Do not use this for ordinary facts or bullet
        summaries.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(data.background_color.replace("#", ""))
        quote_rgb = RGBColor.from_string(data.quote_color.replace("#", ""))
        accent_rgb = RGBColor.from_string(data.accent_color.replace("#", ""))

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._MARGIN_X), Inches(1.15), Inches(0.08), Inches(4.75))
        accent.fill.solid()
        accent.fill.fore_color.rgb = accent_rgb
        accent.line.fill.background()

        mark_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X + 0.35), Inches(1.0), Inches(1.0), Inches(0.8)).text_frame
        mark_frame.text = "\""
        mark_frame.paragraphs[0].font.size = Pt(64)
        mark_frame.paragraphs[0].font.bold = True
        mark_frame.paragraphs[0].font.color.rgb = accent_rgb

        quote_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X + 0.45), Inches(1.85), Inches(self._CONTENT_WIDTH - 0.9), Inches(3.0)).text_frame
        quote_frame.text = data.quote_text
        quote_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        quote_frame.paragraphs[0].font.size = Pt(data.quote_font_size)
        quote_frame.paragraphs[0].font.bold = True
        quote_frame.paragraphs[0].font.color.rgb = quote_rgb
        quote_frame.paragraphs[0].line_spacing = 1.08
        quote_frame.word_wrap = True

        if data.attribution:
            attr_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X + 0.5), Inches(5.2), Inches(self._CONTENT_WIDTH - 1.0), Inches(0.5)).text_frame
            attr_frame.text = data.attribution
            attr_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            attr_frame.paragraphs[0].font.size = Pt(data.attribution_font_size)
            attr_frame.paragraphs[0].font.bold = True
            attr_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(data.attribution_color.replace("#", ""))

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"content", "visual"}))
    def comparison(self, data: ComparisonSlideData):
        """
        Create a side-by-side comparison with two named sections.

        Best for similarities, differences, tradeoffs, pros and cons, old versus
        new, theory versus evidence, or before and after contrasts. Keep each side
        concise.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(data.background_color.replace("#", ""))
        heading_rgb = RGBColor.from_string(data.heading_color.replace("#", ""))
        left_rgb = RGBColor.from_string(data.left_color.replace("#", ""))
        right_rgb = RGBColor.from_string(data.right_color.replace("#", ""))

        title_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X), Inches(0.5), Inches(self._CONTENT_WIDTH), Inches(0.8)).text_frame
        title_frame.text = data.heading
        title_frame.paragraphs[0].font.size = Pt(data.title_font_size)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = heading_rgb
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        column_gap = 0.45
        column_width = (self._CONTENT_WIDTH - column_gap) / 2
        right_left = self._MARGIN_X + column_width + column_gap

        left_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._MARGIN_X), Inches(1.55), Inches(column_width), Inches(5.15))
        left_card.fill.solid()
        left_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        left_card.line.color.rgb = RGBColor(226, 232, 240)
        left_card.shadow.inherit = False
        left_card.shadow.visible = True
        left_card.shadow.blur_radius = Pt(4)
        left_card.shadow.distance = Pt(1)
        left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._MARGIN_X), Inches(1.55), Inches(column_width), Inches(0.08))
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = left_rgb
        left_bar.line.fill.background()
        left_title_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X + 0.35), Inches(1.9), Inches(column_width - 0.7), Inches(0.55)).text_frame
        left_title_frame.text = data.left_title
        left_title_frame.paragraphs[0].font.size = Pt(data.section_title_font_size)
        left_title_frame.paragraphs[0].font.bold = True
        left_title_frame.paragraphs[0].font.color.rgb = left_rgb
        left_title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        left_content_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X + 0.35), Inches(2.65), Inches(column_width - 0.7), Inches(3.65)).text_frame
        left_content_frame.word_wrap = True
        self._render_content(left_content_frame, data.left_content)

        right_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right_left), Inches(1.55), Inches(column_width), Inches(5.15))
        right_card.fill.solid()
        right_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        right_card.line.color.rgb = RGBColor(226, 232, 240)
        right_card.shadow.inherit = False
        right_card.shadow.visible = True
        right_card.shadow.blur_radius = Pt(4)
        right_card.shadow.distance = Pt(1)
        right_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right_left), Inches(1.55), Inches(column_width), Inches(0.08))
        right_bar.fill.solid()
        right_bar.fill.fore_color.rgb = right_rgb
        right_bar.line.fill.background()
        right_title_frame = slide.shapes.add_textbox(Inches(right_left + 0.35), Inches(1.9), Inches(column_width - 0.7), Inches(0.55)).text_frame
        right_title_frame.text = data.right_title
        right_title_frame.paragraphs[0].font.size = Pt(data.section_title_font_size)
        right_title_frame.paragraphs[0].font.bold = True
        right_title_frame.paragraphs[0].font.color.rgb = right_rgb
        right_title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        right_content_frame = slide.shapes.add_textbox(Inches(right_left + 0.35), Inches(2.65), Inches(column_width - 0.7), Inches(3.65)).text_frame
        right_content_frame.word_wrap = True
        self._render_content(right_content_frame, data.right_content)

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"content", "visual"}))
    def timeline(self, data: TimelineSlideData):
        """
        Create a chronological slide of dated or time-ordered events.

        Best for historical periods, milestones, reigns, discoveries, or cause
        chains tied to time. For ordered steps without dates, choose the process
        flow tool instead.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(data.background_color.replace("#", ""))
        heading_rgb = RGBColor.from_string(data.heading_color.replace("#", ""))
        accent_rgb = RGBColor.from_string(data.accent_color.replace("#", ""))

        title_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X), Inches(0.5), Inches(self._CONTENT_WIDTH), Inches(0.8)).text_frame
        title_frame.text = data.heading
        title_frame.paragraphs[0].font.size = Pt(data.title_font_size)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = heading_rgb
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._MARGIN_X), Inches(1.28), Inches(1.45), Inches(0.06))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_rgb
        accent_bar.line.fill.background()

        event_count = max(len(data.events), 1)
        event_height = min(0.85, 4.9 / event_count)
        spacing = (4.9 - event_height) / max(event_count - 1, 1)
        line_left = self._MARGIN_X + 2.45
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(line_left), Inches(1.75), Inches(0.04), Inches(4.8))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(203, 213, 225)
        line.line.fill.background()

        for idx, event in enumerate(data.events):
            top_position = 1.75 + idx * spacing

            node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(line_left - 0.13), Inches(top_position + 0.23), Inches(0.3), Inches(0.3))
            node.fill.solid()
            node.fill.fore_color.rgb = accent_rgb
            node.line.fill.background()
            date_pill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._MARGIN_X), Inches(top_position + 0.12), Inches(2.05), Inches(0.52))
            date_pill.fill.solid()
            date_pill.fill.fore_color.rgb = RGBColor(255, 255, 255)
            date_pill.line.color.rgb = RGBColor(226, 232, 240)
            date_frame = date_pill.text_frame
            date_frame.text = event.date
            date_frame.paragraphs[0].font.size = Pt(min(data.event_font_size, 15))
            date_frame.paragraphs[0].font.bold = True
            date_frame.paragraphs[0].font.color.rgb = accent_rgb
            date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            date_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

            desc_frame = slide.shapes.add_textbox(Inches(line_left + 0.35), Inches(top_position), Inches(self._CONTENT_WIDTH - 2.85), Inches(event_height)).text_frame
            desc_frame.text = event.description
            desc_frame.paragraphs[0].font.size = Pt(min(data.event_font_size, 17))
            desc_frame.paragraphs[0].font.color.rgb = RGBColor(30, 41, 59)
            desc_frame.paragraphs[0].line_spacing = 1.12
            desc_frame.word_wrap = True

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"content", "visual"}, tags={"supports_images"}))
    def image_grid(self, data: ImageGridSlideData):
        """
        Create a grid of related images with short captions.

        Best for visual examples, artifact sets, monuments, species, maps, or
        side-by-side visual inspection. For one main image, choose the body image or
        full image tool.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(data.background_color.replace("#", ""))
        heading_rgb = RGBColor.from_string(data.heading_color.replace("#", ""))
        accent_rgb = RGBColor.from_string(data.accent_color.replace("#", ""))

        title_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X), Inches(0.5), Inches(self._CONTENT_WIDTH), Inches(0.8)).text_frame
        title_frame.text = data.heading
        title_frame.paragraphs[0].font.size = Pt(data.title_font_size)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = heading_rgb
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        num_images = len(data.images)
        if num_images == 0:
            self._publish(slide)
            return slide

        cols = data.grid_columns
        rows = (num_images + cols - 1) // cols
        horizontal_spacing = 0.35
        vertical_spacing = 0.3
        caption_height = 0.48
        img_width = (self._CONTENT_WIDTH - (cols - 1) * horizontal_spacing) / cols
        img_height = (5.25 - rows * caption_height - (rows - 1) * vertical_spacing) / rows

        for idx, image in enumerate(data.images):
            row = idx // cols
            col = idx % cols

            left = self._MARGIN_X + col * (img_width + horizontal_spacing)
            top = 1.65 + row * (img_height + caption_height + vertical_spacing)
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(img_width), Inches(img_height + caption_height + 0.1))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = RGBColor(226, 232, 240)
            card.shadow.inherit = False
            card.shadow.visible = True
            card.shadow.blur_radius = Pt(3)
            card.shadow.distance = Pt(1)

            try:
                self._add_picture_fit(slide, image.image, Inches(left + 0.08), Inches(top + 0.08), Inches(img_width - 0.16), Inches(img_height - 0.08), data.image_fit)
            except Exception:
                self._create_image_placeholder(
                    slide, Inches(left + 0.08), Inches(top + 0.08),
                    Inches(img_width - 0.16), Inches(img_height - 0.08),
                    "[Image]"
                )

            if image.caption:
                caption_frame = slide.shapes.add_textbox(Inches(left + 0.16), Inches(top + img_height + 0.05), Inches(img_width - 0.32), Inches(caption_height)).text_frame
                caption_frame.text = image.caption
                caption_frame.paragraphs[0].font.size = Pt(12)
                caption_frame.paragraphs[0].font.color.rgb = RGBColor(71, 85, 105)
                caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                caption_frame.word_wrap = True

        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._MARGIN_X), Inches(1.28), Inches(1.45), Inches(0.06))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_rgb
        accent_bar.line.fill.background()

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"structural", "visual"}))
    def section_divider(self, data: SectionDividerSlideData):
        """
        Create a major section or chapter transition slide.

        Best when the presentation moves to a new unit, time period, theme, or
        activity block. Do not use this for normal content topics.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(data.background_color.replace("#", ""))
        text_rgb = RGBColor.from_string(data.text_color.replace("#", ""))

        top_rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches((self._BASE_SLIDE_WIDTH_INCHES - 1.5) / 2), Inches(0.9), Inches(1.5), Inches(0.06))
        top_rule.fill.solid()
        top_rule.fill.fore_color.rgb = text_rgb
        top_rule.line.fill.background()

        if data.section_number:
            number_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT), Inches(1.55), Inches(self._CENTER_CONTENT_WIDTH), Inches(0.45)).text_frame
            number_frame.text = data.section_number
            number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            number_frame.paragraphs[0].font.size = Pt(18)
            number_frame.paragraphs[0].font.bold = True
            number_frame.paragraphs[0].font.color.rgb = text_rgb

        title_top = 2.25 if data.section_number else 2.0
        title_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT), Inches(title_top), Inches(self._CENTER_CONTENT_WIDTH), Inches(1.9)).text_frame
        title_frame.text = data.section_title
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(min(data.title_font_size, 60 if len(data.section_title) <= 26 else 52))
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = text_rgb
        title_frame.paragraphs[0].line_spacing = 0.95

        bottom_rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches((self._BASE_SLIDE_WIDTH_INCHES - 2.1) / 2), Inches(title_top + 2.05), Inches(2.1), Inches(0.05))
        bottom_rule.fill.solid()
        bottom_rule.fill.fore_color.rgb = text_rgb
        bottom_rule.fill.transparency = 0.35
        bottom_rule.line.fill.background()

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"content", "structural"}))
    def key_takeaways(self, data: KeyTakeawaysSlideData):
        """
        Create a summary slide of the main points students should remember.

        Best at the end of a section or lesson. Use 3-5 concise takeaway statements
        and avoid introducing new material.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(data.background_color.replace("#", ""))
        heading_rgb = RGBColor.from_string(data.heading_color.replace("#", ""))
        accent_rgb = RGBColor.from_string(data.accent_color.replace("#", ""))

        title_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X), Inches(0.55), Inches(self._CONTENT_WIDTH), Inches(0.75)).text_frame
        title_frame.text = data.heading
        title_frame.paragraphs[0].font.size = Pt(data.title_font_size)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = heading_rgb
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._MARGIN_X), Inches(1.28), Inches(1.45), Inches(0.06))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_rgb
        accent_bar.line.fill.background()

        takeaways = data.takeaways
        use_grid = data.layout == "grid" or (data.layout == "auto" and len(takeaways) <= 4)

        if use_grid:
            cols = 2 if len(takeaways) > 1 else 1
            rows = (len(takeaways) + cols - 1) // cols
            gap_x = 0.35
            gap_y = 0.3
            card_width = (self._CONTENT_WIDTH - (cols - 1) * gap_x) / cols
            card_height = (5.0 - (rows - 1) * gap_y) / rows
            takeaway_font_size = min(data.takeaway_font_size, 18)
            for idx, takeaway in enumerate(takeaways):
                row = idx // cols
                col = idx % cols
                left = self._MARGIN_X + col * (card_width + gap_x)
                top = 1.65 + row * (card_height + gap_y)
                card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(card_width), Inches(card_height))
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(255, 255, 255)
                card.line.color.rgb = RGBColor(226, 232, 240)
                card.line.width = Pt(1)
                pill = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.32), Inches(top + 0.35), Inches(0.46), Inches(0.46))
                pill.fill.solid()
                pill.fill.fore_color.rgb = accent_rgb
                pill.line.fill.background()
                pill.text_frame.text = str(idx + 1) if data.numbered else "•"
                pill.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                pill.text_frame.paragraphs[0].font.size = Pt(13)
                pill.text_frame.paragraphs[0].font.bold = True
                pill.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                takeaway_frame = slide.shapes.add_textbox(Inches(left + 0.35), Inches(top + 1.0), Inches(card_width - 0.7), Inches(card_height - 1.25)).text_frame
                takeaway_frame.text = takeaway
                takeaway_frame.word_wrap = True
                takeaway_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                takeaway_frame.paragraphs[0].font.size = Pt(takeaway_font_size)
                takeaway_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
                takeaway_frame.paragraphs[0].line_spacing = 1.08
            self._publish(slide)
            return slide

        card_gap = 0.16
        card_area_height = 5.15 - (len(takeaways) - 1) * card_gap
        takeaway_font_size = min(data.takeaway_font_size, 16 if len(takeaways) <= 5 else 14)
        chars_per_line = max(32, int((self._CENTER_CONTENT_WIDTH - 1.1) * (130 / takeaway_font_size)))
        height_weights = [1 + (min(3, max(1, (len(takeaway) + chars_per_line - 1) // chars_per_line)) - 1) * 0.45 for takeaway in takeaways]
        weight_total = sum(height_weights)
        card_heights = [max(0.68, card_area_height * weight / weight_total) for weight in height_weights]
        if sum(card_heights) > card_area_height:
            card_heights = [height * card_area_height / sum(card_heights) for height in card_heights]

        top = 1.55
        for idx, takeaway in enumerate(takeaways):
            card_height = card_heights[idx]
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._CENTER_CONTENT_LEFT), Inches(top), Inches(self._CENTER_CONTENT_WIDTH), Inches(card_height))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = RGBColor(226, 232, 240)
            card.line.width = Pt(1)
            pill = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(self._CENTER_CONTENT_LEFT + 0.25), Inches(top + (card_height - 0.42) / 2), Inches(0.42), Inches(0.42))
            pill.fill.solid()
            pill.fill.fore_color.rgb = accent_rgb
            pill.line.fill.background()
            pill_frame = pill.text_frame
            pill_frame.text = str(idx + 1) if data.numbered else "•"
            pill_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            pill_frame.paragraphs[0].font.size = Pt(13)
            pill_frame.paragraphs[0].font.bold = True
            pill_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            takeaway_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT + 0.82), Inches(top + 0.12), Inches(self._CENTER_CONTENT_WIDTH - 1.1), Inches(card_height - 0.24)).text_frame
            takeaway_frame.text = takeaway
            takeaway_frame.word_wrap = True
            takeaway_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            takeaway_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            takeaway_frame.paragraphs[0].font.size = Pt(takeaway_font_size)
            takeaway_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
            takeaway_frame.paragraphs[0].line_spacing = 1.1
            top += card_height + card_gap

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"content", "structural"}))
    def agenda(self, data: AgendaSlideData):
        """
        Create a roadmap slide for the lesson flow or session schedule.

        Best near the beginning when students need to see topics, activities, or
        timing. Do not use this to explain lesson content.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(data.background_color.replace("#", ""))
        heading_rgb = RGBColor.from_string(data.heading_color.replace("#", ""))
        accent_rgb = RGBColor.from_string(data.accent_color.replace("#", ""))

        title_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X), Inches(0.55), Inches(self._CONTENT_WIDTH), Inches(0.8)).text_frame
        title_frame.text = data.heading
        title_frame.paragraphs[0].font.size = Pt(data.title_font_size)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = heading_rgb
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        item_gap = 0.16
        item_height = min(0.78, (5.15 - (len(data.items) - 1) * item_gap) / max(len(data.items), 1))
        start_top = 1.65
        time_width = 2.4

        for idx, agenda in enumerate(data.items):
            top_position = start_top + idx * (item_height + item_gap)
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._CENTER_CONTENT_LEFT), Inches(top_position), Inches(self._CENTER_CONTENT_WIDTH), Inches(item_height))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = RGBColor(226, 232, 240)
            badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(self._CENTER_CONTENT_LEFT + 0.2), Inches(top_position + (item_height - 0.42) / 2), Inches(0.42), Inches(0.42))
            badge.fill.solid()
            badge.fill.fore_color.rgb = accent_rgb
            badge.line.fill.background()
            badge.text_frame.text = str(idx + 1)
            badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            badge.text_frame.paragraphs[0].font.size = Pt(13)
            badge.text_frame.paragraphs[0].font.bold = True
            badge.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

            topic_width = self._CENTER_CONTENT_WIDTH - time_width - 0.9 if agenda.time else self._CENTER_CONTENT_WIDTH - 1.0
            topic_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT + 0.78), Inches(top_position + 0.11), Inches(topic_width), Inches(item_height - 0.22)).text_frame
            topic_frame.text = agenda.topic
            topic_frame.paragraphs[0].font.size = Pt(data.item_font_size)
            topic_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
            topic_frame.word_wrap = True
            topic_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            topic_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            if agenda.time:
                time_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT + self._CENTER_CONTENT_WIDTH - time_width - 0.25), Inches(top_position + 0.12), Inches(time_width), Inches(item_height - 0.24)).text_frame
                time_frame.text = agenda.time
                time_frame.paragraphs[0].font.size = Pt(min(data.item_font_size, 15))
                time_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
                time_frame.paragraphs[0].font.bold = True
                time_frame.paragraphs[0].font.color.rgb = accent_rgb
                time_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"engagement", "visual"}, tags={"supports_images"}))
    def full_image(self, data: FullImageSlideData):
        """
        Create a full-bleed image slide with a short text overlay.

        Best for a strong visual hook: a monument, place, person, artwork, map, or
        scene. Keep text minimal. If the slide needs explanation beside the image,
        choose the body image tool.

        NOTE: It is strongly recommended you use the browse images tool to find a
        suitable image for this slide.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        if data.add_overlay:
            base = Image.open(data.image_path).convert("RGBA")
            c = ImageColor.getrgb(data.overlay_color)
            alpha = int(max(0, min(1, data.overlay_opacity)) * 255)
            overlay = Image.new("RGBA", base.size, (c[0], c[1], c[2], alpha))
            composed = Image.alpha_composite(base, overlay).convert("RGB")
            image = io.BytesIO()
            composed.save(image, format="PNG")
            image.seek(0)
        else:
            image = data.image_path

        try:
            self._add_picture_fit(slide, image, Inches(0), Inches(0), Inches(self._BASE_SLIDE_WIDTH_INCHES), Inches(self._BASE_SLIDE_HEIGHT_INCHES), "cover")
        except Exception:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(200, 200, 200)


        text_rgb = RGBColor.from_string(data.text_color.replace('#', ''))
        title_font_size = min(data.title_font_size, 48 if len(data.title) > 32 else data.title_font_size)
        title_lines = max(1, (len(data.title) + 23) // 24)
        title_height = min(2.45, max(1.0, title_lines * title_font_size / 58))
        subtitle_font_size = min(data.subtitle_font_size, 20)
        subtitle_height = 0
        if data.subtitle:
            subtitle_lines = max(1, (len(data.subtitle) + 68) // 69)
            subtitle_height = min(1.15, max(0.55, subtitle_lines * subtitle_font_size / 62))
        group_height = title_height + subtitle_height + (0.18 if data.subtitle else 0)

        if data.title_position == "top":
            title_top = 0.7
        elif data.title_position == "bottom":
            title_top = self._BASE_SLIDE_HEIGHT_INCHES - group_height - 0.75
        else:
            title_top = (self._BASE_SLIDE_HEIGHT_INCHES - group_height) / 2

        # Title
        title_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT), Inches(title_top), Inches(self._CENTER_CONTENT_WIDTH), Inches(title_height)).text_frame
        title_frame.text = data.title
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(title_font_size)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = text_rgb
        title_frame.paragraphs[0].line_spacing = 0.95
        for run in title_frame.paragraphs[0].runs:
            run.font.color.rgb = text_rgb

        # Subtitle
        if data.subtitle:
            subtitle_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT), Inches(title_top + title_height + 0.18), Inches(self._CENTER_CONTENT_WIDTH), Inches(subtitle_height)).text_frame
            subtitle_frame.text = data.subtitle
            subtitle_frame.word_wrap = True
            subtitle_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            subtitle_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            subtitle_frame.paragraphs[0].font.size = Pt(subtitle_font_size)
            subtitle_frame.paragraphs[0].font.color.rgb = text_rgb
            subtitle_frame.paragraphs[0].line_spacing = 1.05
            for run in subtitle_frame.paragraphs[0].runs:
                run.font.color.rgb = text_rgb

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"content", "visual"}))
    def icon(self, data: IconSlideData):
        """
        Create a compact visual slide for 2-4 concepts represented by icons.

        Best for abstract ideas, grouped examples, roles, features, dynasties,
        values, or quick taxonomy. Keep each item to a short title and one-line
        description.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string((data.background_color or "#F8FAFC").replace("#", ""))
        icon_rgb = RGBColor.from_string(data.icon_color.replace("#", ""))
        accent_rgb = RGBColor.from_string(data.accent_color.replace("#", ""))
        heading_rgb = RGBColor.from_string(data.heading_color.replace("#", ""))

        title_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X), Inches(0.5), Inches(self._CONTENT_WIDTH), Inches(0.8)).text_frame
        title_frame.text = data.heading
        title_frame.paragraphs[0].font.size = Pt(34 if len(data.heading) < 52 else 30)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = heading_rgb
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._MARGIN_X), Inches(1.3), Inches(1.45), Inches(0.06))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_rgb
        accent_bar.line.fill.background()

        # Layout icons
        num_icons = len(data.icons)
        if num_icons == 0:
            self._publish(slide)
            return slide

        if data.layout == "horizontal":
            horizontal_spacing = 0.4
            icon_width = (self._CONTENT_WIDTH - (num_icons - 1) * horizontal_spacing) / num_icons
            for idx, icon in enumerate(data.icons):
                left = self._MARGIN_X + idx * (icon_width + horizontal_spacing)
                self._render_icon_item(slide, icon, left, 2.25, icon_width, 2.35, icon_rgb, accent_rgb)

        elif data.layout == "vertical":
            row_left = self._BASE_SLIDE_WIDTH_INCHES * 0.15
            row_width = self._BASE_SLIDE_WIDTH_INCHES - row_left * 2
            row_height = min(1.05, 4.9 / max(num_icons, 1))
            for idx, icon in enumerate(data.icons):
                self._render_icon_item(slide, icon, row_left, 1.65 + idx * (row_height + 0.18), row_width, row_height, icon_rgb, accent_rgb)

        else:  # grid
            if num_icons <= 4:
                cols = 2
            elif num_icons <= 6:
                cols = 3
            else:
                cols = 3

            rows = (num_icons + cols - 1) // cols
            horizontal_spacing = 0.4
            vertical_spacing = 0.34
            icon_width = (self._CONTENT_WIDTH - (cols - 1) * horizontal_spacing) / cols
            icon_height = min(2.05, (4.95 - (rows - 1) * vertical_spacing) / rows)
            grid_top = 1.65 + max(0, (4.95 - rows * icon_height - (rows - 1) * vertical_spacing) / 2)

            for idx, icon in enumerate(data.icons):
                row = idx // cols
                col = idx % cols
                left = self._MARGIN_X + col * (icon_width + horizontal_spacing)
                top = grid_top + row * (icon_height + vertical_spacing)
                self._render_icon_item(slide, icon, left, top, icon_width, icon_height, icon_rgb, accent_rgb)

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"content", "visual"}))
    def highlight_box(self, data: HighlightBoxSlideData):
        """
        Create a slide of 2-4 colored boxes for standout points.

        Best for key facts, factors, causes, effects, benefits, or short principles
        that deserve equal emphasis. For ordered steps, choose process flow; for two
        opposing sides, choose comparison.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(data.background_color.replace("#", ""))
        heading_rgb = RGBColor.from_string(data.heading_color.replace("#", ""))
        accent_rgb = RGBColor.from_string(data.accent_color.replace("#", ""))

        title_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X), Inches(0.5), Inches(self._CONTENT_WIDTH), Inches(0.8)).text_frame
        title_frame.text = data.heading
        title_frame.paragraphs[0].font.size = Pt(data.title_font_size)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = heading_rgb
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._MARGIN_X), Inches(1.28), Inches(1.45), Inches(0.06))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_rgb
        accent_bar.line.fill.background()

        num_boxes = len(data.boxes)
        if num_boxes == 0:
            self._publish(slide)
            return slide

        if data.layout == "horizontal":
            box_gap = 0.4
            box_width = (self._CONTENT_WIDTH - (num_boxes - 1) * box_gap) / num_boxes
            for idx, box in enumerate(data.boxes):
                self._render_highlight_box(slide, box, Inches(self._MARGIN_X + idx * (box_width + box_gap)), Inches(2.0), Inches(box_width), Inches(4.0))

        elif data.layout == "vertical":
            box_height = (5.5 - (num_boxes + 1) * 0.2) / num_boxes
            for idx, box in enumerate(data.boxes):
                self._render_highlight_box(slide, box, Inches(self._CENTER_CONTENT_LEFT), Inches(1.8 + idx * (box_height + 0.2)), Inches(self._CENTER_CONTENT_WIDTH), Inches(box_height))

        else:  # grid
            cols = 2
            rows = (num_boxes + 1) // 2
            box_gap_x = 0.4
            box_gap_y = 0.3
            box_width = (self._CONTENT_WIDTH - box_gap_x) / 2
            box_height = (5.5 - box_gap_y * (rows - 1)) / rows
            for idx, box in enumerate(data.boxes):
                row = idx // cols
                col = idx % cols
                self._render_highlight_box(slide, box, Inches(self._MARGIN_X + col * (box_width + box_gap_x)), Inches(1.8 + row * (box_height + box_gap_y)), Inches(box_width), Inches(box_height))

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"content", "visual"}))
    def big_number(self, data: BigNumberSlideData):
        """
        Create a high-impact slide around one important number.

        Best for a memorable date, statistic, count, percentage, or scale marker
        that can be explained in one short context line. For multiple dates, choose
        the timeline tool.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string((data.background_color or "#F8FAFC").replace("#", ""))
        number_rgb = RGBColor.from_string(data.number_color.replace("#", ""))
        accent_rgb = RGBColor.from_string(data.accent_color.replace("#", ""))

        panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._CENTER_CONTENT_LEFT), Inches(1.25), Inches(self._CENTER_CONTENT_WIDTH), Inches(4.95))
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
        panel.line.color.rgb = RGBColor(226, 232, 240)
        panel.shadow.inherit = False
        panel.shadow.visible = True
        panel.shadow.blur_radius = Pt(6)
        panel.shadow.distance = Pt(2)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._CENTER_CONTENT_LEFT), Inches(1.25), Inches(self._CENTER_CONTENT_WIDTH), Inches(0.08))
        accent.fill.solid()
        accent.fill.fore_color.rgb = accent_rgb
        accent.line.fill.background()

        number_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT + 0.5), Inches(1.85), Inches(self._CENTER_CONTENT_WIDTH - 1.0), Inches(1.45)).text_frame
        number_frame.text = data.number
        number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        number_frame.paragraphs[0].font.size = Pt(96)
        number_frame.paragraphs[0].font.bold = True
        number_frame.paragraphs[0].font.color.rgb = number_rgb

        context_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT + 0.8), Inches(3.65), Inches(self._CENTER_CONTENT_WIDTH - 1.6), Inches(1.1)).text_frame
        context_frame.text = data.context
        context_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        context_frame.paragraphs[0].font.size = Pt(24 if len(data.context) < 80 else 20)
        context_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(data.context_color.replace("#", ""))
        context_frame.paragraphs[0].line_spacing = 1.1
        context_frame.word_wrap = True

        if data.subtitle:
            subtitle_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT + 1.0), Inches(5.15), Inches(self._CENTER_CONTENT_WIDTH - 2.0), Inches(0.45)).text_frame
            subtitle_frame.text = data.subtitle
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            subtitle_frame.paragraphs[0].font.size = Pt(16)
            subtitle_frame.paragraphs[0].font.bold = True
            subtitle_frame.paragraphs[0].font.color.rgb = accent_rgb
            subtitle_frame.word_wrap = True

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"content", "visual"}))
    def process_flow(self, data: ProcessFlowSlideData):
        """
        Create a step-by-step flow slide.

        Best for procedures, stages, workflows, systems, hierarchies, or chains of
        cause and effect. If the order is primarily chronological with dates, choose
        the timeline tool.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string((data.background_color or "#F8FAFC").replace("#", ""))
        heading_rgb = RGBColor.from_string(data.heading_color.replace("#", ""))
        flow_rgb = RGBColor.from_string(data.flow_color.replace("#", ""))
        accent_rgb = RGBColor.from_string(data.accent_color.replace("#", ""))

        title_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X), Inches(0.5), Inches(self._CONTENT_WIDTH), Inches(0.8)).text_frame
        title_frame.text = data.heading
        title_frame.paragraphs[0].font.size = Pt(data.title_font_size)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = heading_rgb
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._MARGIN_X), Inches(1.28), Inches(1.45), Inches(0.06))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_rgb
        accent_bar.line.fill.background()

        num_steps = len(data.steps)
        if num_steps == 0:
            self._publish(slide)
            return slide

        def render_step(idx: int, step: ProcessStep, left: float, top: float, width: float, height: float):
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = RGBColor(226, 232, 240)
            card.shadow.inherit = False
            card.shadow.visible = True
            card.shadow.blur_radius = Pt(3)
            card.shadow.distance = Pt(1)

            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.25), Inches(top + 0.25), Inches(0.5), Inches(0.5))
            circle.fill.solid()
            circle.fill.fore_color.rgb = flow_rgb
            circle.line.color.rgb = flow_rgb

            num_frame = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.25), Inches(0.5), Inches(0.5)).text_frame
            num_frame.text = str(step.step_number)
            num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            num_frame.paragraphs[0].font.size = Pt(16)
            num_frame.paragraphs[0].font.bold = True
            num_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            num_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

            title_frame = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.95), Inches(width - 0.5), Inches(0.52)).text_frame
            title_frame.text = step.title
            title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            title_frame.paragraphs[0].font.size = Pt(14 if num_steps >= 5 else 15)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(15, 23, 42)
            title_frame.word_wrap = True

            desc_frame = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 1.6), Inches(width - 0.5), Inches(height - 1.85)).text_frame
            desc_frame.text = step.description
            desc_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            desc_frame.paragraphs[0].font.size = Pt(10.5 if num_steps >= 5 else 12)
            desc_frame.paragraphs[0].font.color.rgb = RGBColor(71, 85, 105)
            desc_frame.paragraphs[0].line_spacing = 1.12
            desc_frame.word_wrap = True

        if num_steps >= 5:
            cols = 3
            rows = (num_steps + cols - 1) // cols
            gap_x = 0.35
            gap_y = 0.28
            step_width = (self._CONTENT_WIDTH - (cols - 1) * gap_x) / cols
            step_height = (4.75 - (rows - 1) * gap_y) / rows
            for idx, step in enumerate(data.steps):
                row = idx // cols
                col = idx % cols
                render_step(idx, step, self._MARGIN_X + col * (step_width + gap_x), 1.75 + row * (step_height + gap_y), step_width, step_height)
        else:
            step_gap = 0.3
            step_width = (self._CONTENT_WIDTH - (num_steps - 1) * step_gap) / num_steps
            for idx, step in enumerate(data.steps):
                left = self._MARGIN_X + idx * (step_width + step_gap)
                render_step(idx, step, left, 1.85, step_width, 4.65)
                if idx < num_steps - 1:
                    arrow_width = min(step_gap * 0.75, 0.35)
                    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left + step_width + (step_gap - arrow_width) / 2), Inches(3.92), Inches(arrow_width), Inches(0.12))
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = flow_rgb
                    arrow.line.fill.background()

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"content", "engagement"}))
    def question(self, data: QuestionSlideData):
        """
        Create one focused question for student thinking.

        Best for retrieval practice, prediction, checking understanding, quick
        polls, or a single multiple-choice prompt. For several open-ended prompts,
        choose the discussion tool.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # Background
        if data.background_color:
            background = slide.background
            fill = background.fill
            fill.solid()
            rgb = RGBColor.from_string(data.background_color.replace('#', ''))
            fill.fore_color.rgb = rgb

        question_rgb = RGBColor.from_string(data.question_color.replace('#', ''))

        icon_frame = slide.shapes.add_textbox(Inches((self._BASE_SLIDE_WIDTH_INCHES - 1.2) / 2), Inches(0.55), Inches(1.2), Inches(0.75)).text_frame
        icon_frame.text = data.icon
        icon_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        icon_frame.paragraphs[0].font.size = Pt(46)

        question_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT), Inches(1.5), Inches(self._CENTER_CONTENT_WIDTH), Inches(1.45)).text_frame
        question_frame.text = data.question
        question_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        question_frame.paragraphs[0].font.size = Pt(30 if len(data.question) < 90 else 26)
        question_frame.paragraphs[0].font.bold = True
        question_frame.paragraphs[0].font.color.rgb = question_rgb
        question_frame.word_wrap = True
        question_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        question_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        question_frame.paragraphs[0].line_spacing = 1.1

        if data.options:
            columns = 2 if len(data.options) > 1 else 1
            rows = (len(data.options) + columns - 1) // columns
            option_gap_x = 0.32
            option_gap_y = 0.18
            option_top = 3.35
            option_width = (self._CENTER_CONTENT_WIDTH - (columns - 1) * option_gap_x) / columns
            option_height = min(0.82, (2.35 - (rows - 1) * option_gap_y) / rows)
            for idx, option in enumerate(data.options):
                row = idx // columns
                col = idx % columns
                left = self._CENTER_CONTENT_LEFT + col * (option_width + option_gap_x)
                top = option_top + row * (option_height + option_gap_y)
                option_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(option_width), Inches(option_height))
                option_card.fill.solid()
                option_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
                option_card.fill.transparency = 0.06
                option_card.line.color.rgb = RGBColor(226, 232, 240)
                option_card.line.width = Pt(1)
                badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.18), Inches(top + (option_height - 0.34) / 2), Inches(0.34), Inches(0.34))
                badge.fill.solid()
                badge.fill.fore_color.rgb = question_rgb
                badge.line.fill.background()
                badge.text_frame.text = chr(65 + idx)
                badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                badge.text_frame.paragraphs[0].font.size = Pt(12)
                badge.text_frame.paragraphs[0].font.bold = True
                badge.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                option_frame = slide.shapes.add_textbox(Inches(left + 0.62), Inches(top + 0.1), Inches(option_width - 0.82), Inches(option_height - 0.2)).text_frame
                option_frame.text = option
                option_frame.word_wrap = True
                option_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                option_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                option_frame.paragraphs[0].font.size = Pt(17)
                option_frame.paragraphs[0].font.color.rgb = RGBColor(17, 24, 39)

        if data.hint:
            hint_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(self._CENTER_CONTENT_LEFT + 0.6), Inches(6.18), Inches(self._CENTER_CONTENT_WIDTH - 1.2), Inches(0.55))
            hint_box.fill.solid()
            hint_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            hint_box.fill.transparency = 0.18
            hint_box.line.fill.background()
            hint_frame = hint_box.text_frame
            hint_frame.text = f"💡 Hint: {data.hint}"
            hint_frame.word_wrap = True
            hint_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            hint_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            hint_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            hint_frame.paragraphs[0].font.size = Pt(13)
            hint_frame.paragraphs[0].font.italic = True
            hint_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"content", "engagement"}))
    def discussion(self, data: DiscussionSlideData):
        """
        Create a discussion prompt slide for pair, group, or class conversation.

        Best for interpretation, debate, personal connection, ethical questions, or
        evidence-based reasoning. Use 2-4 prompts and include timing when helpful.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string((data.background_color or "#F8FAFC").replace("#", ""))
        heading_rgb = RGBColor.from_string(data.heading_color.replace("#", ""))
        accent_rgb = RGBColor.from_string(data.accent_color.replace("#", ""))

        icon_chip = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(self._MARGIN_X), Inches(0.52), Inches(0.72), Inches(0.72))
        icon_chip.fill.solid()
        icon_chip.fill.fore_color.rgb = RGBColor(255, 255, 255)
        icon_chip.line.color.rgb = RGBColor(226, 232, 240)
        icon_frame = icon_chip.text_frame
        icon_frame.text = data.icon
        icon_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        icon_frame.paragraphs[0].font.size = Pt(30)

        heading_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X + 0.9), Inches(0.55), Inches(self._CONTENT_WIDTH - 3.3), Inches(0.75)).text_frame
        heading_frame.text = data.heading
        heading_frame.paragraphs[0].font.size = Pt(34)
        heading_frame.paragraphs[0].font.bold = True
        heading_frame.paragraphs[0].font.color.rgb = heading_rgb

        badge_width = 2.67
        badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._BASE_SLIDE_WIDTH_INCHES - self._MARGIN_X - badge_width), Inches(0.72), Inches(badge_width), Inches(0.42))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(255, 255, 255)
        badge.line.color.rgb = RGBColor(226, 232, 240)
        type_frame = badge.text_frame
        type_frame.text = f"📋 {data.discussion_type.replace('_', ' ').title()}"
        type_frame.paragraphs[0].font.size = Pt(14)
        type_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        type_frame.paragraphs[0].font.bold = True
        type_frame.paragraphs[0].font.color.rgb = accent_rgb

        prompt_gap = 0.18
        prompt_height = min(1.05, (4.45 - (len(data.discussion_prompts) - 1) * prompt_gap) / max(len(data.discussion_prompts), 1))
        for idx, prompt in enumerate(data.discussion_prompts, start=1):
            top = 1.65 + (idx - 1) * (prompt_height + prompt_gap)
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._CENTER_CONTENT_LEFT), Inches(top), Inches(self._CENTER_CONTENT_WIDTH), Inches(prompt_height))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = RGBColor(226, 232, 240)
            marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(self._CENTER_CONTENT_LEFT + 0.25), Inches(top + 0.23), Inches(0.42), Inches(0.42))
            marker.fill.solid()
            marker.fill.fore_color.rgb = heading_rgb
            marker.line.fill.background()
            marker.text_frame.text = str(idx)
            marker.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            marker.text_frame.paragraphs[0].font.size = Pt(13)
            marker.text_frame.paragraphs[0].font.bold = True
            marker.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            prompt_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT + 0.9), Inches(top + 0.15), Inches(self._CENTER_CONTENT_WIDTH - 1.15), Inches(prompt_height - 0.3)).text_frame
            prompt_frame.text = prompt
            prompt_frame.word_wrap = True
            prompt_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            prompt_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            prompt_frame.paragraphs[0].font.size = Pt(18)
            prompt_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
            prompt_frame.paragraphs[0].line_spacing = 1.12

        if data.time_allocation:
            time_frame = slide.shapes.add_textbox(Inches((self._BASE_SLIDE_WIDTH_INCHES - 4.0) / 2), Inches(6.55), Inches(4.0), Inches(0.4)).text_frame
            time_frame.text = f"⏱️ {data.time_allocation}"
            time_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            time_frame.paragraphs[0].font.size = Pt(16)
            time_frame.paragraphs[0].font.bold = True
            time_frame.paragraphs[0].font.color.rgb = accent_rgb

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"content", "engagement"}))
    def activity(self, data: ActivitySlideData):
        """
        Create a concrete student task with step-by-step instructions.

        Best for hands-on practice, group work, source analysis, experiments,
        sorting, mapping, or role play. Include time, group size, and materials when
        useful.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string((data.background_color or "#F8FAFC").replace("#", ""))
        title_rgb = RGBColor.from_string(data.title_color.replace("#", ""))
        accent_rgb = RGBColor.from_string(data.accent_color.replace("#", ""))

        icon_chip = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(self._MARGIN_X), Inches(0.52), Inches(0.72), Inches(0.72))
        icon_chip.fill.solid()
        icon_chip.fill.fore_color.rgb = RGBColor(255, 255, 255)
        icon_chip.line.color.rgb = RGBColor(226, 232, 240)
        icon_frame = icon_chip.text_frame
        icon_frame.text = data.icon
        icon_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        icon_frame.paragraphs[0].font.size = Pt(30)

        title_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X + 0.9), Inches(0.55), Inches(self._CONTENT_WIDTH - 0.9), Inches(0.75)).text_frame
        title_frame.text = data.activity_title
        title_frame.paragraphs[0].font.size = Pt(34)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = title_rgb

        metadata_text = []
        if data.time_required:
            metadata_text.append(f"⏱️ {data.time_required}")
        if data.group_size:
            metadata_text.append(f"👥 {data.group_size}")

        if metadata_text:
            meta_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X + 0.9), Inches(1.25), Inches(self._CONTENT_WIDTH - 0.9), Inches(0.4)).text_frame
            meta_frame.text = " | ".join(metadata_text)
            meta_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            meta_frame.paragraphs[0].font.size = Pt(14)
            meta_frame.paragraphs[0].font.bold = True
            meta_frame.paragraphs[0].font.color.rgb = accent_rgb

        instruction_gap = 0.14
        instruction_top = 1.85
        instruction_height = min(0.72, (3.65 - (len(data.instructions) - 1) * instruction_gap) / max(len(data.instructions), 1))
        for idx, instruction in enumerate(data.instructions, start=1):
            top = instruction_top + (idx - 1) * (instruction_height + instruction_gap)
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._CENTER_CONTENT_LEFT), Inches(top), Inches(self._CENTER_CONTENT_WIDTH), Inches(instruction_height))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = RGBColor(226, 232, 240)
            marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(self._CENTER_CONTENT_LEFT + 0.22), Inches(top + (instruction_height - 0.36) / 2), Inches(0.36), Inches(0.36))
            marker.fill.solid()
            marker.fill.fore_color.rgb = title_rgb
            marker.line.fill.background()
            marker.text_frame.text = str(idx)
            marker.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            marker.text_frame.paragraphs[0].font.size = Pt(12)
            marker.text_frame.paragraphs[0].font.bold = True
            marker.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            instruction_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT + 0.75), Inches(top + 0.09), Inches(self._CENTER_CONTENT_WIDTH - 0.95), Inches(instruction_height - 0.18)).text_frame
            instruction_frame.text = instruction
            instruction_frame.word_wrap = True
            instruction_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            instruction_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            instruction_frame.paragraphs[0].font.size = Pt(16)
            instruction_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)

        if data.materials_needed:
            materials_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._CENTER_CONTENT_LEFT), Inches(5.95), Inches(self._CENTER_CONTENT_WIDTH), Inches(0.75))
            materials_box.fill.solid()
            materials_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            materials_box.line.fill.background()
            materials_frame = materials_box.text_frame
            materials_frame.text = f"📦 Materials: {', '.join(data.materials_needed)}"
            materials_frame.paragraphs[0].font.size = Pt(14)
            materials_frame.paragraphs[0].font.bold = True
            materials_frame.paragraphs[0].font.color.rgb = RGBColor(71, 85, 105)
            materials_frame.word_wrap = True

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"content", "engagement"}))
    def real_world_connection(self, data: RealWorldConnectionSlideData):
        """
        Create a relevance slide that connects lesson content to the real world.

        Best for modern applications, careers, civic issues, technology, current
        parallels, or student experience. Provide concrete examples and an optional
        reflection prompt.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string((data.background_color or "#F8FAFC").replace("#", ""))
        heading_rgb = RGBColor.from_string(data.heading_color.replace("#", ""))
        accent_rgb = RGBColor.from_string(data.accent_color.replace("#", ""))

        icon_chip = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(self._MARGIN_X), Inches(0.52), Inches(0.72), Inches(0.72))
        icon_chip.fill.solid()
        icon_chip.fill.fore_color.rgb = RGBColor(255, 255, 255)
        icon_chip.line.color.rgb = RGBColor(226, 232, 240)
        icon_frame = icon_chip.text_frame
        icon_frame.text = data.icon
        icon_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        icon_frame.paragraphs[0].font.size = Pt(30)

        heading_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X + 0.9), Inches(0.55), Inches(self._CONTENT_WIDTH - 0.9), Inches(0.75)).text_frame
        heading_frame.text = data.heading
        heading_frame.paragraphs[0].font.size = Pt(34)
        heading_frame.paragraphs[0].font.bold = True
        heading_frame.paragraphs[0].font.color.rgb = heading_rgb

        connection_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT), Inches(1.45), Inches(self._CENTER_CONTENT_WIDTH), Inches(0.62)).text_frame
        connection_frame.text = data.connection_title
        connection_frame.paragraphs[0].font.size = Pt(24)
        connection_frame.paragraphs[0].font.bold = True
        connection_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        connection_frame.paragraphs[0].font.color.rgb = heading_rgb

        desc_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._CENTER_CONTENT_LEFT), Inches(2.2), Inches(self._CENTER_CONTENT_WIDTH), Inches(1.25))
        desc_box.fill.solid()
        desc_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        desc_box.line.color.rgb = RGBColor(226, 232, 240)
        desc_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT + 0.35), Inches(2.42), Inches(self._CENTER_CONTENT_WIDTH - 0.7), Inches(0.85)).text_frame
        desc_frame.text = data.description
        desc_frame.paragraphs[0].font.size = Pt(16)
        desc_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
        desc_frame.word_wrap = True
        desc_frame.paragraphs[0].line_spacing = 1.15

        cols = 2 if len(data.examples) > 1 else 1
        rows = (len(data.examples) + cols - 1) // cols
        card_gap_x = 0.25
        card_gap_y = 0.16
        card_width = (self._CENTER_CONTENT_WIDTH - (cols - 1) * card_gap_x) / cols
        card_height = min(0.72, (1.75 - (rows - 1) * card_gap_y) / max(rows, 1))
        for idx, example in enumerate(data.examples):
            row = idx // cols
            col = idx % cols
            left = self._CENTER_CONTENT_LEFT + col * (card_width + card_gap_x)
            top = 3.75 + row * (card_height + card_gap_y)
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(card_width), Inches(card_height))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = RGBColor(226, 232, 240)
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.18), Inches(top + (card_height - 0.16) / 2), Inches(0.16), Inches(0.16))
            dot.fill.solid()
            dot.fill.fore_color.rgb = accent_rgb
            dot.line.fill.background()
            example_frame = slide.shapes.add_textbox(Inches(left + 0.48), Inches(top + 0.11), Inches(card_width - 0.65), Inches(card_height - 0.22)).text_frame
            example_frame.text = example
            example_frame.word_wrap = True
            example_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            example_frame.paragraphs[0].font.size = Pt(14)
            example_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)

        if data.call_to_action:
            cta_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._CENTER_CONTENT_LEFT), Inches(6.15), Inches(self._CENTER_CONTENT_WIDTH), Inches(0.58))
            cta_box.fill.solid()
            cta_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            cta_box.line.fill.background()
            cta_frame = cta_box.text_frame
            cta_frame.text = f"💭 {data.call_to_action}"
            cta_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cta_frame.paragraphs[0].font.size = Pt(16)
            cta_frame.paragraphs[0].font.bold = True
            cta_frame.paragraphs[0].font.color.rgb = accent_rgb
            cta_frame.word_wrap = True

        self._publish(slide)
        return slide


    @slide(SlideMeta(categories={"content", "engagement"}))
    def story(self, data: StorySlideData):
        """
        Create a narrative slide built around a short story or case.

        Best for opening a topic, humanizing abstract content, creating curiosity,
        or making a key idea memorable. Keep it brief and factual; do not use this
        for lists of facts.
        """
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string((data.background_color or "#F8FAFC").replace("#", ""))
        title_rgb = RGBColor.from_string(data.title_color.replace("#", ""))
        accent_rgb = RGBColor.from_string(data.accent_color.replace("#", ""))

        icon_chip = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(self._MARGIN_X), Inches(0.52), Inches(0.72), Inches(0.72))
        icon_chip.fill.solid()
        icon_chip.fill.fore_color.rgb = RGBColor(255, 255, 255)
        icon_chip.line.color.rgb = RGBColor(226, 232, 240)
        corner_frame = icon_chip.text_frame
        corner_frame.text = data.icon
        corner_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        corner_frame.paragraphs[0].font.size = Pt(30)

        badge_width = 2.67
        type_badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._BASE_SLIDE_WIDTH_INCHES - self._MARGIN_X - badge_width), Inches(0.72), Inches(badge_width), Inches(0.42))
        type_badge.fill.solid()
        type_badge.fill.fore_color.rgb = RGBColor(255, 255, 255)
        type_badge.line.color.rgb = RGBColor(226, 232, 240)
        type_frame = type_badge.text_frame
        type_frame.text = data.story_type.replace('_', ' ').title()
        type_frame.paragraphs[0].font.size = Pt(13)
        type_frame.paragraphs[0].font.bold = True
        type_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        type_frame.paragraphs[0].font.color.rgb = accent_rgb

        title_frame = slide.shapes.add_textbox(Inches(self._MARGIN_X + 0.9), Inches(0.55), Inches(self._CONTENT_WIDTH - 3.6), Inches(0.75)).text_frame
        title_frame.text = data.story_title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        title_frame.paragraphs[0].font.color.rgb = title_rgb

        story_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._CENTER_CONTENT_LEFT), Inches(1.7), Inches(self._CENTER_CONTENT_WIDTH), Inches(4.2))
        story_box.fill.solid()
        story_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        story_box.line.color.rgb = RGBColor(226, 232, 240)
        story_box.shadow.inherit = False
        story_box.shadow.visible = True
        story_box.shadow.blur_radius = Pt(4)
        story_box.shadow.distance = Pt(1)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._CENTER_CONTENT_LEFT), Inches(1.7), Inches(0.08), Inches(4.2))
        accent.fill.solid()
        accent.fill.fore_color.rgb = accent_rgb
        accent.line.fill.background()

        story_frame = slide.shapes.add_textbox(Inches(self._CENTER_CONTENT_LEFT + 0.45), Inches(2.1), Inches(self._CENTER_CONTENT_WIDTH - 0.9), Inches(3.35)).text_frame
        story_frame.text = data.story_content
        story_frame.paragraphs[0].font.size = Pt(16)
        story_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
        story_frame.word_wrap = True
        story_frame.paragraphs[0].line_spacing = 1.35
        story_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        if data.moral_or_lesson:
            lesson_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self._CENTER_CONTENT_LEFT), Inches(6.15), Inches(self._CENTER_CONTENT_WIDTH), Inches(0.58))
            lesson_box.fill.solid()
            lesson_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            lesson_box.line.fill.background()
            lesson_frame = lesson_box.text_frame
            lesson_frame.text = f"💡 {data.moral_or_lesson}"
            lesson_frame.paragraphs[0].font.size = Pt(16)
            lesson_frame.paragraphs[0].font.bold = True
            lesson_frame.paragraphs[0].font.color.rgb = accent_rgb
            lesson_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            lesson_frame.word_wrap = True

        self._publish(slide)
        return slide


    def _render_icon_item(self, slide, icon: IconItem, left, top, width, height, icon_rgb, accent_rgb):
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(226, 232, 240)
        card.line.width = Pt(1)
        card.shadow.inherit = False
        card.shadow.visible = True
        card.shadow.blur_radius = Pt(4)
        card.shadow.distance = Pt(1)

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.08), Inches(height))
        accent.fill.solid()
        accent.fill.fore_color.rgb = accent_rgb
        accent.line.fill.background()

        icon_size = min(0.72, max(0.52, height * 0.34))
        icon_left = left + 0.34
        icon_top = top + (height - icon_size) / 2
        chip = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(icon_left), Inches(icon_top), Inches(icon_size), Inches(icon_size))
        chip.fill.solid()
        chip.fill.fore_color.rgb = icon_rgb
        chip.fill.transparency = 0.88
        chip.line.fill.background()

        icon_frame = slide.shapes.add_textbox(Inches(icon_left), Inches(icon_top + 0.04), Inches(icon_size), Inches(icon_size - 0.08)).text_frame
        icon_frame.text = icon.icon_text
        icon_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        icon_frame.paragraphs[0].font.size = Pt(27 if height < 1.4 else 32)
        icon_frame.paragraphs[0].font.color.rgb = icon_rgb

        text_left = left + 1.25
        text_width = width - 1.55
        title_frame = slide.shapes.add_textbox(Inches(text_left), Inches(top + 0.34), Inches(text_width), Inches(0.42)).text_frame
        title_frame.text = icon.title
        title_frame.paragraphs[0].font.size = Pt(17)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(15, 23, 42)
        title_frame.word_wrap = True
        title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        desc_frame = slide.shapes.add_textbox(Inches(text_left), Inches(top + 0.86), Inches(text_width), Inches(max(height - 1.08, 0.36))).text_frame
        desc_frame.text = icon.description
        desc_frame.paragraphs[0].font.size = Pt(13)
        desc_frame.paragraphs[0].font.color.rgb = RGBColor(71, 85, 105)
        desc_frame.paragraphs[0].line_spacing = 1.12
        desc_frame.word_wrap = True
        desc_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


    def _fit_text(self, text_frame: TextFrame):
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


    def _scale(self, slide: Slide):
        width_ratio = self.prs.slide_width / self._BASE_SLIDE_WIDTH
        height_ratio = self.prs.slide_height / self._BASE_SLIDE_HEIGHT

        if width_ratio == 1 and height_ratio == 1:
            return

        for shape in slide.shapes:
            shape.left = int(shape.left * width_ratio)
            shape.top = int(shape.top * height_ratio)
            shape.width = int(shape.width * width_ratio)
            shape.height = int(shape.height * height_ratio)


    def _publish(self, slide: Slide):
        self._scale(slide)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            self._fit_text(shape.text_frame)


    def _render_highlight_box(self, slide, box: HighlightBox, left, top, width, height):
        rgb = RGBColor.from_string(box.color.replace('#', ''))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = RGBColor(226, 232, 240)
        shape.line.width = Pt(1)
        shape.shadow.inherit = False
        shape.shadow.visible = True
        shape.shadow.blur_radius = Pt(4)
        shape.shadow.distance = Pt(1)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
        accent.fill.solid()
        accent.fill.fore_color.rgb = rgb
        accent.line.fill.background()

        title_frame = slide.shapes.add_textbox(left + Inches(0.35), top + Inches(0.38), width - Inches(0.7), Inches(0.55)).text_frame
        title_frame.text = box.title
        title_frame.paragraphs[0].font.size = Pt(18)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = rgb
        title_frame.word_wrap = True

        content_frame = slide.shapes.add_textbox(left + Inches(0.35), top + Inches(1.02), width - Inches(0.7), height - Inches(1.2)).text_frame
        content_frame.text = box.content
        content_frame.paragraphs[0].font.size = Pt(15)
        content_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
        content_frame.paragraphs[0].line_spacing = 1.15
        content_frame.word_wrap = True


    def _image_aspect(self, image) -> float | None:
        try:
            if isinstance(image, io.BytesIO):
                pos = image.tell()
                image.seek(0)
                with Image.open(image) as img:
                    width, height = img.size
                image.seek(pos)
            else:
                with Image.open(image) as img:
                    width, height = img.size
            return width / height if height else None
        except Exception:
            return None


    def _add_picture_fit(self, slide, image, left, top, width, height, fit: Literal["contain", "cover", "stretch"]):
        if fit == "stretch":
            return slide.shapes.add_picture(image, left, top, width=width, height=height)

        aspect = self._image_aspect(image)
        if not aspect:
            return slide.shapes.add_picture(image, left, top, width=width, height=height)

        box_aspect = width / height
        if fit == "contain":
            if aspect >= box_aspect:
                image_width = width
                image_height = int(width / aspect)
            else:
                image_height = height
                image_width = int(height * aspect)
            return slide.shapes.add_picture(image, left + int((width - image_width) / 2), top + int((height - image_height) / 2), width=image_width, height=image_height)

        picture = slide.shapes.add_picture(image, left, top, width=width, height=height)
        if aspect > box_aspect:
            crop = (1 - box_aspect / aspect) / 2
            picture.crop_left = crop
            picture.crop_right = crop
        else:
            crop = (1 - aspect / box_aspect) / 2
            picture.crop_top = crop
            picture.crop_bottom = crop
        return picture


    def _create_image_placeholder(self, slide, left, top, width, height, text):
        image_placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        image_placeholder.fill.solid()
        image_placeholder.fill.fore_color.rgb = RGBColor(226, 237, 248)
        image_placeholder.line.color.rgb = RGBColor(148, 163, 184)
        image_placeholder.line.width = Pt(1)

        text_frame = image_placeholder.text_frame
        text_frame.text = "Image unavailable" if not text or text.startswith("Error loading image:") else text
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(15)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(71, 85, 105)

        return image_placeholder


    def _render_content(self, text_frame: TextFrame, entry: ContentType):
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        if text_frame.paragraphs and not any(r.text for r in text_frame.paragraphs[0].runs):
            text_frame.clear()

        def clean(v: str) -> str: return re.sub(r'\\u[0-9a-fA-F]{4}', "", v)

        def apply_bullet(p):
            from pptx.oxml.xmlchemy import OxmlElement
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
            for el in list(pPr):
                if el.tag in {qn('a:buAutoNum'), qn('a:buChar'), qn('a:buNone')}:
                    pPr.remove(el)
            buChar = OxmlElement('a:buChar')
            buChar.set('char', '\u2022')
            pPr.insert(0, buChar)
            p.level = 0
            p.space_after = Pt(12)  # Increased from 6 to 12 for better spacing
            p.space_before = Pt(6)  # Increased from 2 to 6
            p.line_spacing = 1.4  # Increased from 1.2 to 1.4

            # Add proper indentation for bullet
            from pptx.util import Inches
            p.left_indent = Inches(0.3)  # Indent bullet from left
            p.first_line_indent = Inches(-0.25)  # Hanging indent for bullet

        def apply_autonum(p, start_at: int = 1, num_type: str = 'arabicPeriod'):
            from pptx.oxml.xmlchemy import OxmlElement
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
            for el in list(pPr):
                if el.tag in {qn('a:buAutoNum'), qn('a:buChar'), qn('a:buNone')}:
                    pPr.remove(el)
            buAutoNum = OxmlElement('a:buAutoNum')
            buAutoNum.set('type', num_type)
            buAutoNum.set('startAt', str(start_at))
            pPr.insert(0, buAutoNum)
            p.level = 0
            p.space_after = Pt(12)  # Increased from 6 to 12
            p.space_before = Pt(6)  # Increased from 2 to 6
            p.line_spacing = 1.4  # Increased from 1.2 to 1.4

            # Add proper indentation for numbers
            from pptx.util import Inches
            p.left_indent = Inches(0.3)  # Indent number from left
            p.first_line_indent = Inches(-0.25)  # Hanging indent for number

        text_frame.word_wrap = True

        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        stack = [entry]
        while stack:
            content = stack.pop(0)
            if isinstance(content, list):
                for c in content:
                    stack.append(c)
                continue

            if isinstance(content, str):
                p = text_frame.add_paragraph()
                p.text = clean(content)
                p.font.size = Pt(16)
                p.space_after = Pt(8)
                continue

            if isinstance(content, TextBlock):
                p = text_frame.add_paragraph()
                p.text = clean(content.value)
                if content.formatting:
                    f = content.formatting
                    run = p.runs[0]
                    run.font.bold = f.bold
                    run.font.italic = f.italic
                    run.font.underline = f.underline
                    if f.font_color:
                        rgb = RGBColor.from_string(f.font_color.replace('#', ''))
                        run.font.color.rgb = rgb
                    if f.font_size:
                        run.font.size = Pt(f.font_size)
                    else:
                        run.font.size = Pt(16)
                    if f.alignment:
                        from pptx.enum.text import PP_ALIGN
                        align_map = {
                            "left": PP_ALIGN.LEFT,
                            "center": PP_ALIGN.CENTER,
                            "right": PP_ALIGN.RIGHT,
                            "justify": PP_ALIGN.JUSTIFY,
                        }
                        p.alignment = align_map.get(f.alignment.lower(), PP_ALIGN.LEFT)
                    p.level = f.bullet_level
                else:
                    p.font.size = Pt(16)
                p.space_after = Pt(8)
                continue

            if isinstance(content, ListBlock):
                if content.items:
                    for idx, item in enumerate(content.items, start=1):
                        p = text_frame.add_paragraph()
                        p.text = clean(item)
                        p.font.size = Pt(14)  # Slightly smaller for lists
                        if content.numbered:
                            apply_autonum(p, start_at=idx)
                        else:
                            apply_bullet(p)
                continue

            raise ValueError("Unexpected value encountered:" + str(content))


def get_metadata(exclude: set[str] = {
    Templates.welcome.__name__,
    Templates.body_with_video.__name__,
    Templates.thank_you.__name__,
}) -> str:
    out = []
    for name in dir(Templates):
        fn = getattr(Templates, name)
        if not callable(fn) or fn.__name__ not in SLIDE_META or fn.__name__ in exclude:
            continue

        meta = SLIDE_META[fn.__name__]
        labels = sorted(meta.categories) + sorted(meta.tags)
        label_text = f" ({', '.join(labels)})" if labels else ""

        doc = (inspect.getdoc(fn) or "").replace("\n", " ").replace("  ", " ").strip()
        out.append("- `%s`%s: _%s_" % (fn.__name__, label_text, doc))

    return "\n".join(out).strip()
