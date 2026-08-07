import asyncio
from dataclasses import dataclass, field
from datetime import datetime, timezone
import inspect
from io import BytesIO
import json
from typing import Annotated, Any, Literal
from aiofiles.threadpool.binary import AsyncBufferedReader

from app.models.presentation import CaptionerResponse, ImageSearchResults, NextSlideSpec, PresentationOutline, ShikshaAgentEvent, ShikshaCheckpointEvent, SlideSpec, WelcomeSlideInfo, YouTubeVideoResult
from app.utils.storage import Storage
from pptx import presentation
from pptx.util import Length

from pydantic import Field, HttpUrl, PositiveInt
from pydantic_ai import Agent, AgentRunResultEvent, AgentStreamEvent, BinaryContent, FunctionToolCallEvent, FunctionToolResultEvent, FunctionToolset, ModelRetry, PartEndEvent, RunContext, TextContent, TextPart, ThinkingPart, ToolsetTool, WrapperToolset
from pydantic_ai.capabilities import Thinking
from pydantic_ai.usage import UsageLimits

from app.services.presentation import docparser, template, utils
from app.config import DESIGNER_BODY_SLIDE_PROMPT, DESIGNER_FIRST_SLIDE_PROMPT, FINALIZER_ADD_SLIDE_PROMPT, FINALIZER_REVIEW_PROMPT, FINALIZER_BROWSE_PROMPT, PLANNER_USER_PROMPT, settings, CAPTIONER_SYSTEM_PROMPT, PLANNER_SYSTEM_PROMPT, DESIGNER_SYSTEM_PROMPT, FINALIZER_SYSTEM_PROMPT


@dataclass
class _ImageResolverMixin:
    _known_urls: dict[str, str] = field(default_factory=dict, init=False)

    def map_url(self, url: str) -> str:
        if url not in self._known_urls:
            self._known_urls[url] = utils.randomize_url(url, len(self._known_urls))
        return self._known_urls[url]

    def resolve_url(self, url: str) -> str:
        real = next((k for k, v in self._known_urls.items() if v == url), None)
        if real is None: raise ModelRetry("404 Not Found: %s" % url)
        return real


@dataclass
class DesignerDeps(_ImageResolverMixin):
    storage: Storage
    prs: presentation.Presentation
    templates: template.Templates
    figures_dir: str
    outline: PresentationOutline
    slide: SlideSpec | None
    metadata: dict[str, Any]


@dataclass
class FinalizerDeps(_ImageResolverMixin):
    prs: presentation.Presentation
    templates: template.Templates
    metadata: dict[str, Any]


def _getdoc(fn: object) -> str: return (inspect.getdoc(fn) or "").replace("\n", " ").replace("  ", " ").strip()


async def resolve_image(ctx: _ImageResolverMixin, image: str, storage: Storage | None, local_dir: str | None) -> str | BytesIO:
    if is_url := image.startswith("http"):
        image = ctx.resolve_url(image)
    try:
        return await utils.resolve_image(image, is_url, storage, local_dir)
    except Exception as e:
        raise ModelRetry(str(e)) from e


designer_toolset = FunctionToolset[DesignerDeps](max_retries=3)
@designer_toolset.tool(description=_getdoc(template.Templates.welcome), metadata={"action": "Creating welcome slide"})
def ppt_add_welcome_slide(ctx: RunContext[DesignerDeps], data: template.WelcomeSlideData) -> WelcomeSlideInfo:
    ctx.deps.templates.welcome(data)
    assert isinstance(ctx.deps.templates.prs.slide_width, Length) and isinstance(ctx.deps.templates.prs.slide_height, Length)
    return WelcomeSlideInfo(slide_width_inches=ctx.deps.templates.prs.slide_width.inches, slide_height_inches=ctx.deps.templates.prs.slide_height.inches)

@designer_toolset.tool(description=_getdoc(template.Templates.body), metadata={"action": "Adding body slide"})
def ppt_add_body_slide(ctx: RunContext[DesignerDeps], data: template.BodySlideData) -> str:
    ctx.deps.templates.body(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.body_with_image), metadata={"action": "Adding body slide with image"})
async def ppt_add_body_with_image_slide(ctx: RunContext[DesignerDeps], data: template.BodySlideWithImageData) -> str:
    data.image_path = await resolve_image(ctx.deps, data.image_path, ctx.deps.storage, ctx.deps.figures_dir)
    ctx.deps.templates.body_with_image(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.two_column), metadata={"action": "Adding 2-column slide"})
def ppt_add_two_column_slide(ctx: RunContext[DesignerDeps], data: template.TwoColumnSlideData) -> str:
    ctx.deps.templates.two_column(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.quote), metadata={"action": "Adding quote slide"})
def ppt_add_quote_slide(ctx: RunContext[DesignerDeps], data: template.QuoteSlideData) -> str:
    ctx.deps.templates.quote(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.comparison), metadata={"action": "Adding comparison slide"})
def ppt_add_comparison_slide(ctx: RunContext[DesignerDeps], data: template.ComparisonSlideData) -> str:
    ctx.deps.templates.comparison(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.timeline), metadata={"action": "Adding timeline slide"})
def ppt_add_timeline_slide(ctx: RunContext[DesignerDeps], data: template.TimelineSlideData) -> str:
    ctx.deps.templates.timeline(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.image_grid), metadata={"action": "Adding image grid slide"})
async def ppt_add_image_grid_slide(ctx: RunContext[DesignerDeps], data: template.ImageGridSlideData) -> str:
    results = await asyncio.gather(*[resolve_image(ctx.deps, img.image, ctx.deps.storage, ctx.deps.figures_dir) for img in data.images])
    for img, resolved in zip(data.images, results):
        img.image = resolved
    ctx.deps.templates.image_grid(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.section_divider), metadata={"action": "Adding section divider slide"})
def ppt_add_section_divider_slide(ctx: RunContext[DesignerDeps], data: template.SectionDividerSlideData) -> str:
    ctx.deps.templates.section_divider(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.key_takeaways), metadata={"action": "Adding key takeaways slide"})
def ppt_add_key_takeaways_slide(ctx: RunContext[DesignerDeps], data: template.KeyTakeawaysSlideData) -> str:
    ctx.deps.templates.key_takeaways(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.agenda), metadata={"action": "Adding agenda slide"})
def ppt_add_agenda_slide(ctx: RunContext[DesignerDeps], data: template.AgendaSlideData) -> str:
    ctx.deps.templates.agenda(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.full_image), metadata={"action": "Adding full image slide"})
async def ppt_add_full_image_slide(ctx: RunContext[DesignerDeps], data: template.FullImageSlideData) -> str:
    data.image_path = await resolve_image(ctx.deps, data.image_path, ctx.deps.storage, ctx.deps.figures_dir)
    ctx.deps.templates.full_image(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.icon), metadata={"action": "Adding icon slide"})
def ppt_add_icon_slide(ctx: RunContext[DesignerDeps], data: template.IconSlideData) -> str:
    ctx.deps.templates.icon(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.highlight_box), metadata={"action": "Adding highlight box slide"})
def ppt_add_highlight_box_slide(ctx: RunContext[DesignerDeps], data: template.HighlightBoxSlideData) -> str:
    ctx.deps.templates.highlight_box(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.big_number), metadata={"action": "Adding big number slide"})
def ppt_add_big_number_slide(ctx: RunContext[DesignerDeps], data: template.BigNumberSlideData) -> str:
    ctx.deps.templates.big_number(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.process_flow), metadata={"action": "Adding process flow slide"})
def ppt_add_process_flow_slide(ctx: RunContext[DesignerDeps], data: template.ProcessFlowSlideData) -> str:
    ctx.deps.templates.process_flow(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.question), metadata={"action": "Adding question slide"})
def ppt_add_question_slide(ctx: RunContext[DesignerDeps], data: template.QuestionSlideData) -> str:
    ctx.deps.templates.question(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.discussion), metadata={"action": "Adding discussion slide"})
def ppt_add_discussion_slide(ctx: RunContext[DesignerDeps], data: template.DiscussionSlideData) -> str:
    ctx.deps.templates.discussion(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.activity), metadata={"action": "Adding activity slide"})
def ppt_add_activity_slide(ctx: RunContext[DesignerDeps], data: template.ActivitySlideData) -> str:
    ctx.deps.templates.activity(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.real_world_connection), metadata={"action": "Adding real-world connection slide"})
def ppt_add_real_world_connection_slide(ctx: RunContext[DesignerDeps], data: template.RealWorldConnectionSlideData) -> str:
    ctx.deps.templates.real_world_connection(data)
    return "Slide created successfully."

@designer_toolset.tool(description=_getdoc(template.Templates.story), metadata={"action": "Adding story slide"})
def ppt_add_story_slide(ctx: RunContext[DesignerDeps], data: template.StorySlideData) -> str:
    ctx.deps.templates.story(data)
    return "Slide created successfully."

@designer_toolset.tool(metadata={"action": "Browsing images"})
async def browse_images(
    ctx: RunContext[DesignerDeps],
    queries: Annotated[list[str], Field(description="List of queries for image search. These are inserted directly into `gsrsearch`, so keep it short and search-engine-like.", min_length=1, max_length=5, examples=[
        ["plant cell", "mitochondrion"], ["Pont du Gard", "Apollo lunar module"], ["Apis mellifera"], ["teectonic plates map", "heart anatomy"], ["Hokusai wave"]
    ])],
    max_results_per_query: PositiveInt = 12,
    mime: Annotated[Literal["image/png", "image/jpg", "image/gif", "image/svg+xml"] | None, Field(description="Leave as null to return images of any MIME")] = None
) -> ImageSearchResults:
    """
    Search Wikimedia Commons files using MediaWiki keyword search.

    Query guidance:
    - Use keywords, not sentences.
    - Usually use 1-4 words.
    - Put the main subject first.
    - Add only one clarifier and only when needed: location, species, object type, view, or format.
    - Prefer multiple narrow queries over one long query.
    - Do not add "image", "picture", "photo", "Wikimedia", "Commons", "free", or "reusable".
    - Avoid broad single terms unless the subject is already unique.
    - Set `mime="image/gif"` to gather animated images (great for engagement)

    Bad queries:
    - "an educational image of a plant cell"
    - "free reusable Wikimedia photo of the Pont du Gard"
    - "a clear diagram showing tectonic plates around the world"
    - "Apollo lunar module"

    Pass returned image URLs to image fields and cite attribution from extmetadata in captions.
    """

    if not ctx.deps.slide or "supports_images" not in template.SLIDE_META[ctx.deps.slide.slide_type].tags:
        raise ModelRetry("You cannot browse for images while working on a `%s` slide." % (ctx.deps.slide.slide_type if ctx.deps.slide else "welcome"))

    response = await utils.wikimedia_image_search(queries, max_results_per_query, mime)
    if response.message == utils.MSG_NO_RESULTS_FOUND:
        raise ModelRetry("No results. Retry with 1-4 keyword queries only: subject first, no sentences, no image/photo/free/Wikimedia words.")

    for res in response.results:
        res.url = ctx.deps.map_url(res.url)
    return response


@designer_toolset.tool(metadata={"action": "Reading next slide spec"})
def next_slide_spec(ctx: RunContext[DesignerDeps]) -> NextSlideSpec | None:
    """
    Advance to the next slide to work on and return that slide, or null of no more slides are left.
    """
    slide_counter = 1
    for section in ctx.deps.outline.sections:
        for slide in section.slides:
            if slide_counter < len(ctx.deps.metadata["slides_completed"]):
                slide_counter += 1
                continue
            ctx.deps.slide = slide
            return NextSlideSpec(section_title=section.section_title, slide=slide, slide_number=slide_counter)
    ctx.deps.slide = None
    return None


class SlideTrackerToolset(WrapperToolset[DesignerDeps]):

    async def call_tool(self, name: str, tool_args: dict[str, Any], ctx: RunContext[DesignerDeps], tool: ToolsetTool[DesignerDeps]) -> Any:
        if len(ctx.deps.metadata["slides_completed"]) > ctx.deps.outline.total_slides:
            raise ModelRetry("You have created all slides, you cannot invoke any more tools.")

        spec = ctx.deps.slide
        slide_id = str(ctx.deps.outline.compute_slide_id(spec) if spec is not None else -1)  # work around some dumb shit about mongodb not supporting int keys - raises bson.errors.InvalidDocument 'documents must have only string keys, key was -1'
        is_slide_maker = tool.toolset == designer_toolset and tool.tool_def.name.startswith("ppt_add_")
        if is_slide_maker:
            expected_tool = "ppt_add_%s_slide" % (spec.slide_type if spec else "welcome")
            if tool.tool_def.name != expected_tool:
                raise ModelRetry("You cannot use `%s` tool to build a `%s` slide, use `%s` tool instead." % (tool.tool_def.name, spec.slide_type if spec else "welcome", expected_tool))
            if slide_id in ctx.deps.metadata["slide_ids_created"]:
                raise ModelRetry("This slide has already been created - you cannot invoke this tool anymore. You are likely forgetting to call the `%s` tool." % next_slide_spec.__name__)

        result = await super().call_tool(name, tool_args, ctx, tool)

        if is_slide_maker:
            if spec is not None and "supports_images" in template.SLIDE_META[spec.slide_type].tags:
                ctx.deps.metadata["images_used"] = ctx.deps.metadata.get("images_used", 0) + 1
            if spec is not None and "engagement" in template.SLIDE_META[spec.slide_type].categories:
                ctx.deps.metadata["engagement_slides"] = ctx.deps.metadata.get("engagement_slides", 0) + 1
            t = datetime.now(timezone.utc).isoformat()
            ctx.deps.metadata["slide_ids_created"][slide_id] = t
            ctx.deps.metadata["slides_completed"].append(t)
            ctx.deps.metadata["slides_created"] = len(ctx.deps.prs.slides)
            ctx.deps.metadata["slide_types_used"].append(spec.slide_type if spec else "welcome")
        return result


captioner = Agent(model=settings.pres_captioner, name="captioner", output_type=CaptionerResponse, retries=3, system_prompt=CAPTIONER_SYSTEM_PROMPT.safe_substitute())
planner = Agent(model=settings.pres_planner, name="planner", retries=3, output_type=PresentationOutline, system_prompt=PLANNER_SYSTEM_PROMPT.safe_substitute(template_metadata=template.get_metadata()), capabilities=[Thinking()])
designer = Agent(model=settings.pres_designer, name="designer", deps_type=DesignerDeps, retries=3, system_prompt=DESIGNER_SYSTEM_PROMPT.safe_substitute())
finalizer = Agent(model=settings.pres_finalizer, name="finalizer", deps_type=FinalizerDeps, retries=3, system_prompt=FINALIZER_SYSTEM_PROMPT.safe_substitute())


@finalizer.tool(metadata={"action": "Browsing videos"})
async def find_videos(ctx: RunContext[FinalizerDeps], query: str, max_results: int = 10) -> list[YouTubeVideoResult]:
    """ Search for videos on YouTube. """
    try:
        response = await utils.youtube_video_search(query, max_results)
    except ValueError as e:
        raise ModelRetry(str(e)) from e

    for res in response:
        res.url = HttpUrl(ctx.deps.map_url(str(res.url)))
        res.thumbnail_url = HttpUrl(ctx.deps.map_url(str(res.thumbnail_url)))
    return response

@finalizer.tool(description=_getdoc(template.Templates.body_with_video), metadata={"action": "Adding video slide"})
async def ppt_add_body_video_slide(ctx: RunContext[FinalizerDeps], data: template.BodySlideWithVideoData) -> str:
    data.video = HttpUrl(ctx.deps.resolve_url(str(data.video)))
    data.thumbnail_image = await resolve_image(ctx.deps, data.thumbnail_image, None, None)
    ctx.deps.templates.body_with_video(data)
    return "Slide created successfully."

@finalizer.tool(description=_getdoc(template.Templates.thank_you), metadata={"action": "Adding thank you slide"})
def ppt_add_thank_you_slide(ctx: RunContext[FinalizerDeps], data: template.ThankYouSlideData) -> str:
    ctx.deps.templates.thank_you(data)
    return "Slide created successfully."

def _adapt_event(event: AgentStreamEvent) -> ShikshaAgentEvent | None:
    if isinstance(event, FunctionToolCallEvent):
        args = event.part.args
        return ShikshaAgentEvent(type="ToolCallRequestEvent", content=[{
            "id": event.part.tool_call_id,
            "name": event.part.tool_name,
            "arguments": args if isinstance(args, str) else json.dumps(args or {}),
        }])

    if isinstance(event, FunctionToolResultEvent):
        return ShikshaAgentEvent(type="ToolCallExecutionEvent", content=[{
            "call_id": event.result.tool_call_id,
            "name": event.result.tool_name,
            "is_error": getattr(event.result, "outcome", None) != "success",
            "content": event.result.content,
        }])

    if isinstance(event, PartEndEvent):
        if isinstance(event.part, ThinkingPart) and event.part.has_content():
            return ShikshaAgentEvent(type="Thinking", content=event.part.content)
        if isinstance(event.part, TextPart) and event.part.has_content():
            return ShikshaAgentEvent(type="TextMessage", content=event.part.content)
    return None


async def plan(textbook_path: str, textbook_mime: str, data: AsyncBufferedReader, figures: list[docparser.FigureInfo], slides: int | None, metadata: dict[str, Any], instruction: str | None):
    if "outline" in metadata:
        return

    yield ShikshaCheckpointEvent(message="Reading document to build presentation spec...")

    await data.seek(0)
    if textbook_mime.startswith("text/"):
        f = TextContent((await data.read()).decode(), metadata={"source": textbook_path})
    else:
        f = BinaryContent(await data.read(), media_type=textbook_mime)

    outline = None
    async with planner.run_stream_events([f, PLANNER_USER_PROMPT.safe_substitute(
        slides=slides or "auto",
        local_images="\n".join("- " + f.model_dump_json() for f in figures),
        instruction=instruction or "none"
    )], usage_limits=UsageLimits(request_limit=100)) as events:
        async for raw_event in events:
            if isinstance(raw_event, AgentRunResultEvent):
                outline = raw_event.result.output
                break
            event = _adapt_event(raw_event)
            if event is not None:
                yield event

    if outline is None:
        raise RuntimeError("Failed to generate a presentation outline.")

    outline_metadata = outline.model_dump()
    outline_metadata["total_slides"] = outline.total_slides
    yield ShikshaCheckpointEvent(
        message=f"Presentation outline created - {outline.total_slides} slides planned across {len(outline.sections)} sections",
        metadata={"outline": outline_metadata},
        reason="op"
    )


async def design(storage: Storage, prs: presentation.Presentation, data: AsyncBufferedReader, figures_dir: str, outline: PresentationOutline, metadata: dict[str, Any], instruction: str | None):
    # Initialize metadata
    if "slides_completed" not in metadata: metadata["slides_completed"] = []
    if "slides_created" not in metadata: metadata["slides_created"] = 0
    if "slide_ids_created" not in metadata: metadata["slide_ids_created"] = {}
    if "images_used" not in metadata: metadata["images_used"] = 0
    if "slide_types_used" not in metadata: metadata["slide_types_used"] = []
    if "engagement_slides" not in metadata: metadata["engagement_slides"] = 0

    # Add presentation outline information
    subs = dict(
        outline_title=outline.title,
        outline_total_slides=outline.total_slides,
        outline_sections=len(outline.sections),
        learning_objectives="\n".join("- " + obj for obj in outline.learning_objectives),
        instruction=instruction or "none",
        local_images="\n".join("- " + f.model_dump_json() for f in await docparser.read_figures(storage, data, figures_dir)),
        next_tool=next_slide_spec.__name__
    )

    yield ShikshaCheckpointEvent(message="Executing slide creation based on plan", metadata=metadata)

    templates = template.Templates(prs)
    deps = DesignerDeps(storage, prs, templates, figures_dir, outline, None, metadata)
    # Create welcome slide
    if metadata["slides_created"] == 0:
        yield ShikshaCheckpointEvent(message="Creating engaging welcome slide")
        slide_count = len(prs.slides)
        async with designer.run_stream_events(DESIGNER_FIRST_SLIDE_PROMPT.safe_substitute(**subs), usage_limits=UsageLimits(request_limit=8), deps=deps, toolsets=[
            SlideTrackerToolset(designer_toolset.filtered(lambda _, tool_def: tool_def.name == ppt_add_welcome_slide.__name__))
        ]) as events:
            async for raw_event in events:
                if isinstance(raw_event, AgentRunResultEvent):
                    break
                event = _adapt_event(raw_event)
                if event is not None:
                    yield event
        if len(prs.slides) <= slide_count:
            raise RuntimeError("No slide was created")
        yield ShikshaCheckpointEvent(metadata=metadata, reason="op")

    # Process each section and its slides from the outline
    while len(metadata["slides_completed"]) <= deps.outline.total_slides:
        task = DESIGNER_BODY_SLIDE_PROMPT.safe_substitute(**subs)
        async with designer.run_stream_events(task, usage_limits=UsageLimits(request_limit=64), deps=deps, toolsets=[SlideTrackerToolset(designer_toolset)]) as events:
            async for raw_event in events:
                if isinstance(raw_event, AgentRunResultEvent):
                    continue
                event = _adapt_event(raw_event)
                if event is None:
                    continue
                yield event
                if deps.slide is not None:
                    message = "Creating slide %d/%d: %s (%s)" % (len(metadata["slides_completed"]), deps.outline.total_slides, deps.slide.title, deps.slide.slide_type)
                    yield ShikshaCheckpointEvent(message=message, metadata=metadata, reason="op")
        yield ShikshaCheckpointEvent(metadata=metadata, reason="op")

    # Final engagement quality report
    total_slides = len(prs.slides)
    image_percentage = (metadata["images_used"] / total_slides) * 100 if total_slides > 0 else 0
    unique_types = len(set(metadata["slide_types_used"]))
    engagement_percentage = (metadata["engagement_slides"] / total_slides) * 100 if total_slides > 0 else 0

    yield ShikshaCheckpointEvent(
        message=f"Slides complete! Images: {image_percentage:.0f}% | Variety: {unique_types} types | Engagement: {engagement_percentage:.0f}%",
        metadata=metadata
    )


async def finalize(prs: presentation.Presentation, metadata: dict[str, Any]):
    slide_content = json.dumps(docparser.list_slide_content(prs))

    yield ShikshaCheckpointEvent(message="Finalizing slide content")
    deps = FinalizerDeps(prs, template.Templates(prs), metadata)
    videos = metadata["videos"] if "videos" in metadata else []
    if len(videos) == 0:
        yield ShikshaCheckpointEvent(message="Browsing YouTube for videos")
        async with finalizer.run_stream_events(FINALIZER_BROWSE_PROMPT.safe_substitute(
            slide_content=slide_content,
            videos_found=json.dumps([]),
            videos_relevant=json.dumps([])
        ), usage_limits=UsageLimits(request_limit=8), deps=deps) as events:
            async for raw_event in events:
                if isinstance(raw_event, AgentRunResultEvent):
                    break
                event = _adapt_event(raw_event)
                if event is None:
                    continue
                yield event
                if event.type == "ToolCallExecutionEvent":
                    videos.append(event.model_dump_json())
        if len(videos) > 0:
            metadata["videos"] = videos
            yield ShikshaCheckpointEvent(metadata=metadata, reason="op")

    if len(videos) > 0 and "relevant_videos" not in metadata:
        yield ShikshaCheckpointEvent(message="Reviewing collected videos for a video slide")
        relevant_videos = []
        async with finalizer.run_stream_events(FINALIZER_REVIEW_PROMPT.safe_substitute(
            slide_content=slide_content,
            videos_found=json.dumps(metadata["videos"]),
            videos_relevant=json.dumps([])
        ), usage_limits=UsageLimits(request_limit=8), deps=deps) as events:
            async for raw_event in events:
                if isinstance(raw_event, AgentRunResultEvent):
                    break
                event = _adapt_event(raw_event)
                if event is None:
                    continue
                yield event
                if event.type == "ToolCallExecutionEvent":
                    relevant_videos.append(event.model_dump_json())
        if len(relevant_videos) > 0:
            metadata["relevant_videos"] = relevant_videos
            yield ShikshaCheckpointEvent(metadata=metadata, reason="op")

    if "thank_you_slide_created" not in metadata:
        yield ShikshaCheckpointEvent(message="Creating a thank-you slide")
        slide_count = len(prs.slides)
        async with finalizer.run_stream_events(FINALIZER_ADD_SLIDE_PROMPT.safe_substitute(
            slide_content=slide_content,
            videos_found=json.dumps(metadata.get("videos", [])),
            videos_relevant=json.dumps(metadata.get("relevant_videos", []))
        ), usage_limits=UsageLimits(request_limit=8), deps=deps) as events:
            async for raw_event in events:
                if isinstance(raw_event, AgentRunResultEvent):
                    break
                event = _adapt_event(raw_event)
                if event is not None:
                    yield event
        if len(prs.slides) <= slide_count:
            raise RuntimeError("No slide was created")
        metadata["thank_you_slide_created"] = datetime.now().isoformat()
        yield ShikshaCheckpointEvent(metadata=metadata, reason="op")
