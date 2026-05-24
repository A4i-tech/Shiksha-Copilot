import asyncio
import io
import json, logging, pathlib
from typing import IO, Any, Iterator
import zlib

from PIL import Image
from app.models.presentation import FigureInfo
from app.utils.storage import Storage
from pptx import presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from pydantic_ai import AgentRunError, BinaryContent
import pymupdf4llm
import pymupdf


async def _caption(data: IO[bytes], figure_path: str, page_text: str, sem: asyncio.Semaphore) -> FigureInfo | None:
    from app.services.presentation.agent import captioner
    async with sem:
        try:
            data.seek(0)
            resp = await captioner.run(["Page text:\n" + page_text, BinaryContent(data.read(), media_type="image/png")])
        except AgentRunError:
            return None
        if resp.output.caption is None or resp.output.context is None:
            return None
        return FigureInfo(path=figure_path, caption=resp.output.caption, context=resp.output.context)


def _page_figures(doc: pymupdf.Document, page: pymupdf.Page) -> Iterator[tuple[int, int, str, bytes]]:
    result = []
    for data in page.get_images():
        image = doc.extract_image(data[0])
        if image["width"] > page.rect.width * 0.85 or image["height"] > page.rect.height * 0.85:
            continue
        yield image["width"], image["height"], image["ext"], image["image"]

    buf = []
    for drawing in page.get_drawings():
        rect = drawing["rect"]
        if rect.width < 64 or rect.height < 64:
            continue

        # skip if rect too big
        if rect.width > page.rect.width * 0.85 or rect.height > page.rect.height * 0.85 or rect.get_area() > page.rect.get_area() * 0.7:
            continue

        # skip if current rect is already contained in an existing rect
        if any(r.x0 <= rect.x0 and r.y0 <= rect.y0 and r.x1 >= rect.x1 and r.y1 >= rect.y1 for r in buf):
            continue

        # replace previous rects contained in current rect
        buf = [r for r in buf if not (rect.x0 <= r.x0 and rect.y0 <= r.y0 and rect.x1 >= r.x1 and rect.y1 >= r.y1)]
        buf.append(rect)

    for rect in buf:
        pix = page.get_pixmap(clip=rect, matrix=pymupdf.Matrix(2, 2), alpha=False)
        yield int(rect.width), int(rect.height), "png", pix.tobytes("png")

    return page, result


async def _actually_read_figures(storage: Storage, textbook: IO[bytes], out_dir: str, caption_concurrency: int) -> list[FigureInfo]:
    logger = logging.getLogger(__name__)
    sem = asyncio.Semaphore(caption_concurrency)
    tasks: list[asyncio.Task[FigureInfo | None]] = []

    async def proc(page_text: str, figure_name: str, content: bytes):
        storage_path = storage.path(out_dir, figure_name)
        if await storage.exists(storage_path):
            return None
        data = io.BytesIO()
        Image.open(io.BytesIO(content)).convert("RGB").save(data, format="PNG")
        await storage.write_bytes(storage_path, data.getvalue())
        return await _caption(data, figure_name, page_text, sem)

    textbook.seek(0)
    try:
        doc = pymupdf.open(stream=textbook)
    except pymupdf.FileDataError:
        return []

    try:
        for page in doc.pages():
            assert isinstance(page, pymupdf.Page)
            page_text = None
            for w, h, ext, content in _page_figures(doc, page):
                if w < 180 or h < 120:
                    continue
                figure_name = '%d.png' % zlib.crc32(content)
                if page_text is None: page_text = page.get_text()
                tasks.append(asyncio.create_task(proc(page_text, figure_name, content)))
    finally:
        doc.close()

    logger.info("Captioning %d figures", len(tasks))
    figures = await asyncio.gather(*tasks)
    return [f for f in figures if f is not None]


async def read_figures(storage: Storage, textbook: IO[bytes], out_dir: str, caption_concurrency: int = 5) -> list[FigureInfo]:
    storage_mf_path = storage.path(out_dir, "manifest.json")
    if await storage.exists(storage_mf_path):
        return list(map(FigureInfo.model_validate, json.loads(await storage.read_text(storage_mf_path))))

    figures = await _actually_read_figures(storage, textbook, out_dir, caption_concurrency)
    await storage.write_text(storage_mf_path, json.dumps(list(map(lambda v: v.model_dump(), figures))))
    return figures


def _transform(content: IO[bytes], mime: str) -> str | None:
    if mime == "application/pdf":
        content.seek(0)
        with pymupdf.open(stream=content, filetype="pdf") as doc:
            return pymupdf4llm.to_markdown(doc)
    return None


async def transform(storage: Storage, content: IO[bytes], mime: str, out_path: str) -> str | None:
    storage_path = out_path + ".md"
    if not await storage.exists(storage_path):
        markdown = await asyncio.to_thread(_transform, content, mime)
        if markdown is None:
            return None
        await storage.write_text(storage_path, markdown)
    return pathlib.Path(storage_path).name


def list_slide_content(prs: presentation.Presentation) -> list[dict[str, Any]]:
    slides = []
    for s in prs.slides:
        title = ""
        text = []
        for shp in s.shapes:
            if shp.is_placeholder and shp.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                title = shp.text.strip()
            if hasattr(shp, "text") and shp.text:
                text.append(shp.text)
        slides.append({"title": title, "text": " ".join(text)})
    return slides


async def save_pptx(storage: Storage, prs: presentation.Presentation, path: str):
    data = io.BytesIO()
    prs.save(data)
    await storage.write_bytes(path, data.getvalue())
